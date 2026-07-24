#!/usr/bin/env python3
"""Run the declared, reproducible cross-format corpus through production router."""
import argparse, hashlib, json, re, shutil, subprocess, sys, tempfile
from collections import Counter, defaultdict
from pathlib import Path
sys.path.insert(0, str(Path(__file__).parent))
from cross_format_regression import normalize_bundle, fingerprint, write_diff, semantically_equal
from regression_environment import environment_manifest

ROOT=Path(__file__).resolve().parents[1]; MANIFEST=ROOT/'tests/cross_format/cases.json'; KNOWN={"docx_basic","docx_horizontal_merge","docx_vertical_merge","xlsx_merged","xlsx_two_blocks","xlsx_formula","pptx_text","pptx_table","pptx_picture","pptx_grouped","pdf_digital","pdf_ocr_text","pdf_ocr_table_accepted","pdf_ocr_table_rejected","csv_big5","json_unicode","html_complex","txt_native","md_native","csv_gb18030"}; FORMATS={"docx","xlsx","pptx","pdf","csv","json","html","txt","md"}
SKILL_VERSION=(ROOT/'VERSION').read_text(encoding='utf-8').strip()
OCR_FONT_CANDIDATES=(
    ROOT/'tests'/'fixtures'/'fonts'/'DejaVuSans-Bold.ttf',
    Path('/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf'),
    Path('/Library/Fonts/Arial Bold.ttf'),
    Path('C:/Windows/Fonts/arialbd.ttf'),
)
def resolve_ocr_font(image_font, size=72, candidates=OCR_FONT_CANDIDATES):
    """Choose a readable TrueType font, with Pillow's scalable built-in fallback."""
    for candidate in candidates:
        if candidate.is_file():
            return image_font.truetype(str(candidate), size)
    return image_font.load_default(size=size)
def load_cases(path=MANIFEST):
    data=json.loads(Path(path).read_text(encoding='utf-8')); seen=set()
    for case in data['cases']:
        assert case['case_id'] not in seen, 'duplicate case_id'; seen.add(case['case_id'])
        assert case['format'] in FORMATS and case['profile'] in {'core','optional-pandoc'}
        assert case['fixture'] in KNOWN and case['expected_status']
        for key in ('table_count','asset_count'):
            value=case[key]; assert value.get('min',0)>=0 and value.get('max', value.get('min',0)) >= value.get('min',0)
        assert case['max_chunk_chars'] <= 2000 and all(len(x)==len(set(x)) for x in (case['required_warning_codes'],case['allowed_warning_codes'],case['forbidden_warning_codes']))
        assert case['determinism'] == {'normalized_reruns_must_match':True}
        patterns=case.get('required_list_markdown_patterns',[])
        assert isinstance(patterns,list) and all(isinstance(pattern,str) and pattern for pattern in patterns)
        for pattern in patterns:
            try: re.compile(pattern)
            except re.error as exc: raise AssertionError(f'invalid required_list_markdown_patterns regex: {exc}')
        ancestors=case.get('required_element_ancestor_types',{})
        assert isinstance(ancestors,dict) and all(isinstance(key,str) and isinstance(value,list) and value and all(isinstance(ancestor,str) and ancestor for ancestor in value) for key,value in ancestors.items())
        ocr=case.get('required_ocr',{})
        assert isinstance(ocr,dict) and set(ocr) <= {'content_patterns','element_types','locator_fields','min_accepted_tables','min_rejected_tables','candidate_reason_codes'}
        assert all(isinstance(ocr.get(key,[]),list) and all(isinstance(value,str) and value for value in ocr.get(key,[])) for key in ('content_patterns','element_types','locator_fields','candidate_reason_codes'))
        for pattern in ocr.get('content_patterns',[]):
            try: re.compile(pattern)
            except re.error as exc: raise AssertionError(f'invalid required_ocr.content_patterns regex: {exc}')
        assert all(isinstance(ocr.get(key,0),int) and ocr.get(key,0) >= 0 for key in ('min_accepted_tables','min_rejected_tables'))
    return data['cases']
def load_workflow_cases(path=MANIFEST):
    data=json.loads(Path(path).read_text(encoding='utf-8')); seen=set()
    for case in data.get('workflow_cases',[]):
        assert case['case_id'] not in seen, 'duplicate workflow case_id'; seen.add(case['case_id'])
        assert case['profile'] in {'core','optional-pandoc'} and case['workflow'] in {'readable_projection_host_review', 'explicit_user_request_review', 'ocr_advisory_target_review'}
        assert case['fixture'] in KNOWN and case['expected_status'] == 'passed'
        assert case['expected_request_status'] == 'generated' and case['expected_ai_review_status'] == 'not_provided'
        assert case['expected_readable_projection_status'] == 'deterministic_only'
        assert case['determinism'] == {'normalized_reruns_must_match':True}
        assert all(isinstance(pattern,str) and pattern for pattern in case['required_projection_patterns'])
    return data.get('workflow_cases',[])
def generate(name, directory):
    directory.mkdir(parents=True,exist_ok=True)
    if name=='html_complex': return ROOT/'tests/fixtures/html/complex_tables.html'
    if name.startswith('docx'):
        from docx import Document
        d=Document(); p=d.add_paragraph(); p.add_run('bold').bold=True; p.add_run(' italic').italic=True
        if name!='docx_basic':
            t=d.add_table(rows=2,cols=2); a=t.cell(0,0); a.merge(t.cell(0,1) if name.endswith('horizontal_merge') else t.cell(1,0)); a.text='Merged'; t.cell(1,0).text='A';t.cell(1,1).text='B'
        p=directory/(name+'.docx');d.save(p);return p
    if name.startswith('xlsx'):
        import openpyxl
        w=openpyxl.Workbook();s=w.active;s['A1']='Header';s['B1']='Value';s['A2']='A';s['B2']=1
        if name=='xlsx_merged':s.merge_cells('A1:B1')
        if name=='xlsx_two_blocks':s['A5']='Second';s['B5']='Value';s['A6']='B';s['B6']=2
        if name=='xlsx_formula':s['C3']='=A2+B2'
        p=directory/(name+'.xlsx');w.save(p);return p
    if name.startswith('pptx'):
        from pptx import Presentation
        from pptx.util import Inches
        p=Presentation(); slide=p.slides.add_slide(p.slide_layouts[6])
        if name=='pptx_text':
            from pptx.oxml.ns import qn
            from lxml import etree
            box=slide.shapes.add_textbox(Inches(1),Inches(1),Inches(4),Inches(2)); frame=box.text_frame
            first=frame.paragraphs[0]; first.text='Explicit bullet'; etree.SubElement(first._p.get_or_add_pPr(),qn('a:buChar'),{'char':'•'})
            second=frame.add_paragraph(); second.text='Explicit numbered bullet'; etree.SubElement(second._p.get_or_add_pPr(),qn('a:buAutoNum'),{'type':'arabicPeriod'})
        elif name=='pptx_picture':
            from PIL import Image
            image=directory/'pixel.png';Image.new('RGB',(4,4),'red').save(image);slide.shapes.add_picture(str(image),Inches(1),Inches(1),Inches(1),Inches(1))
        else:
            shape=slide.shapes.add_table(2,2,Inches(1),Inches(1),Inches(4),Inches(2)); shape.table.cell(0,0).text='Merged';shape.table.cell(0,0).merge(shape.table.cell(0,1));shape.table.cell(1,0).text='A';shape.table.cell(1,1).text='B'
            if name=='pptx_grouped':
                # Preserve the actual table XML under a group shape, matching
                # the production converter's recursive group traversal path.
                slide.shapes._spTree.remove(shape._element); group=slide.shapes.add_group_shape([]); group._element.append(shape._element)
        out=directory/(name+'.pptx');p.save(out);return out
    if name=='pdf_digital':
        from reportlab.pdfgen import canvas
        p=directory/'digital.pdf';c=canvas.Canvas(str(p));c.drawString(100,700,'Digital PDF Test 123');c.save();return p
    if name.startswith('pdf_ocr_'):
        from PIL import Image, ImageDraw, ImageFont
        font=resolve_ocr_font(ImageFont, 72)
        image=Image.new('RGB',(1800,1100),'white'); draw=ImageDraw.Draw(image)
        if name=='pdf_ocr_text':
            draw.text((130,180),'OCR SAMPLE 12345',font=font,fill='black')
            draw.text((130,300),'SCANNED DOCUMENT',font=font,fill='black')
        elif name=='pdf_ocr_table_accepted':
            rows=[('ITEM','QUANTITY'),('MOTOR','TWO'),('PUMP','FOUR'),('VALVE','TWELVE')]
            x0,y0,w,h=120,100,1550,220
            for row,(left,right) in enumerate(rows):
                y=y0+row*h; draw.rectangle((x0,y,x0+w,y+h),outline='black',width=5); draw.line((x0+900,y,x0+900,y+h),fill='black',width=5)
                draw.text((x0+35,y+62),left,font=font,fill='black'); draw.text((x0+980,y+62),right,font=font,fill='black')
        else:
            draw.text((130,180),'MODEL ABC-100',font=font,fill='black')
            draw.text((130,330),'VOLTAGE 480 V',font=font,fill='black')
        p=directory/(name+'.pdf'); image.save(p,'PDF',resolution=144.0); return p
    if name=='csv_big5':
        p=directory/'big5.csv';p.write_bytes('名稱,數量\n鋼管,10\n'.encode('big5'));return p
    if name=='csv_gb18030':
        p=directory/'gb18030.csv';p.write_bytes('张伟,北京市,GB18030\n李娜,上海市,100\n'.encode('gb18030'));return p
    if name=='json_unicode':
        p=directory/'unicode.json';p.write_text(json.dumps({'名稱':'繁體中文','items':['α','資料']},ensure_ascii=False),encoding='utf-8');return p
    if name=='txt_native':
        p=directory/'sample.txt';p.write_text('Title: Native Text Route\nLine 1: Plain text content\nLine 2: System log entry\n',encoding='utf-8');return p
    if name=='md_native':
        p=directory/'sample.md';p.write_text('# Native Markdown Route\n\n- Item 1\n- Item 2\n\nParagraph text here.\n',encoding='utf-8');return p
def warning_codes(report): return {w.get('code') for w in report.get('warnings',[]) if w.get('code')}
def _safe_table_asset(tables_dir, value):
    """Resolve canonical table assets without permitting bundle traversal."""
    if not isinstance(value, str): return False
    candidate=(tables_dir/value).resolve()
    try: candidate.relative_to(tables_dir.resolve())
    except ValueError: return False
    return candidate.is_file()
def _table_index_errors(index_path, table_ids):
    try: entries=json.loads(index_path.read_text(encoding='utf-8'))
    except (OSError, json.JSONDecodeError): return ['TABLE_INDEX_MALFORMED']
    # Production table_export writes a list. Reject other shapes structurally,
    # rather than accidentally iterating mapping keys or crashing the runner.
    if not isinstance(entries,list): return ['TABLE_INDEX_MALFORMED']
    errors=[]
    for entry in entries:
        if not isinstance(entry,dict): errors.append('TABLE_INDEX_MALFORMED'); continue
        if entry.get('id') not in table_ids: errors.append('REFERENCE_ERROR'); continue
        assets=entry.get('assets',{})
        if not isinstance(assets,dict) or any(not _safe_table_asset(index_path.parent,path) for path in assets.values()): errors.append('REFERENCE_ERROR')
    return errors
def _load_workflow_json(path, missing_code, malformed_code):
    """Load a workflow object without allowing artifact corruption to abort a run."""
    try:
        value=json.loads(path.read_text(encoding='utf-8'))
    except FileNotFoundError:
        return None, [missing_code]
    except json.JSONDecodeError:
        return None, [malformed_code]
    except OSError:
        return None, ['WORKFLOW_ARTIFACT_UNREADABLE']
    if not isinstance(value, dict):
        return None, [malformed_code]
    return value, []
def assert_contract(case,bundle,router_rc):
    errors=[]
    if not (bundle/'conversion-report.json').exists() or not (bundle/'document.json').exists(): return ['ROUTER_BUNDLE_MISSING'], {}
    report=json.loads((bundle/'conversion-report.json').read_text(encoding='utf-8')); doc=json.loads((bundle/'document.json').read_text(encoding='utf-8'))
    if router_rc or report.get('status') not in case['expected_status']: errors.append('STATUS_MISMATCH')
    if report.get('bundle_validation',{}).get('status') != case['expected_bundle_validation']: errors.append('BUNDLE_VALIDATION_FAILED')
    from validate_bundle import validate_bundle
    if validate_bundle(str(bundle)).get('status') != case['expected_bundle_validation']: errors.append('INDEPENDENT_BUNDLE_VALIDATION_FAILED')
    if case.get('expected_engine') and case['expected_engine'] not in str(report.get('engine','')): errors.append('ENGINE_MISMATCH')
    ocr=case.get('required_ocr',{})
    if ocr:
        details=report.get('details',{})
        candidates=details.get('ocr_table_candidates',[])
        assessment=details.get('ocr_table_assessment',{})
        if not details.get('ocr_used') or not candidates: errors.append('OCR_EVIDENCE_MISSING')
        text='\n'.join(str(element.get('content','')) for element in doc.get('elements',[]))
        if any(not re.search(pattern,text) for pattern in ocr.get('content_patterns',[])): errors.append('OCR_CONTENT_MISSING')
        types={e.get('type') for e in doc.get('elements',[])}
        if not set(ocr.get('element_types',[])) <= types: errors.append('OCR_EVIDENCE_MISSING')
        ocr_elements=[e for e in doc.get('elements',[]) if 'ocr' in str(e.get('engine','')).lower() or 'tesseract' in str(e.get('engine','')).lower()]
        if not ocr_elements or any(not all(field in e.get('source_locator',{}) for field in ocr.get('locator_fields',[])) for e in ocr_elements): errors.append('OCR_EVIDENCE_MISSING')
        if assessment.get('accepted_count',0) < ocr.get('min_accepted_tables',0): errors.append('OCR_TABLE_ACCEPTANCE_MISSING')
        if assessment.get('rejected_count',0) < ocr.get('min_rejected_tables',0): errors.append('OCR_TABLE_REJECTION_MISSING')
        candidate_codes={code for candidate in candidates for code in candidate.get('reason_codes',[])}
        if not set(ocr.get('candidate_reason_codes',[])) <= candidate_codes: errors.append('OCR_TABLE_REJECTION_MISSING')
    codes=warning_codes(report); required=set(case['required_warning_codes']); allowed=required|set(case['allowed_warning_codes'])
    if not required <= codes: errors.append('REQUIRED_WARNING_MISSING')
    if codes-allowed: errors.append('UNEXPECTED_WARNING')
    if codes & set(case['forbidden_warning_codes']): errors.append('FORBIDDEN_WARNING')
    types={e.get('type') for e in doc.get('elements',[])}
    if not set(case['required_element_types']) <= types: errors.append('ELEMENT_TYPE_MISSING')
    elements=doc.get('elements',[]); ids=[e.get('id') for e in elements]
    if len(ids)!=len(set(ids)) or any(not ident for ident in ids): errors.append('ELEMENT_ID_INVALID')
    idset=set(ids)
    by_id={element.get('id'): element for element in elements}
    for pattern in case.get('required_list_markdown_patterns',[]):
        if not any(element.get('type') == 'list' and re.search(pattern, element.get('content',''), re.MULTILINE) for element in elements): errors.append('LIST_SEMANTIC_MISSING')
    for element_type, required_ancestors in case.get('required_element_ancestor_types',{}).items():
        matches=[element for element in elements if element.get('type') == element_type]
        def ancestor_types(element):
            result=[]; visited={element.get('id')}; parent=by_id.get(element.get('parent_id')); cycle=False
            while parent:
                parent_id=parent.get('id')
                if parent_id in visited:
                    cycle=True; break
                visited.add(parent_id); result.append(parent.get('type')); parent=by_id.get(parent.get('parent_id'))
            return result, cycle
        candidates=[ancestor_types(element) for element in matches]
        if any(cycle for _,cycle in candidates): errors.append('REFERENCE_ERROR')
        if not any(not cycle and all(kind in types for kind in required_ancestors) for types,cycle in candidates): errors.append('ANCESTOR_SEMANTIC_MISSING')
    for element in elements:
        locator=element.get('source_locator',{})
        if not all(field in locator for field in case['required_locator_fields']): errors.append('LOCATOR_MISSING')
        if element.get('parent_id') and element['parent_id'] not in idset: errors.append('REFERENCE_ERROR')
        if any(child not in idset for child in element.get('children',element.get('child_ids',[]))): errors.append('REFERENCE_ERROR')
    tables=list((bundle/'tables').glob('*.json')) if (bundle/'tables').exists() else []; tables=[p for p in tables if p.name!='index.json']
    table_ids={path.stem for path in tables}
    index_path=bundle/'tables'/'index.json'
    if index_path.exists():
        errors.extend(_table_index_errors(index_path,table_ids))
    lo,hi=case['table_count'].get('min',0),case['table_count'].get('max',10**9)
    if not lo<=len(tables)<=hi: errors.append('TABLE_COUNT')
    assets=list((bundle/'assets').rglob('*')) if (bundle/'assets').exists() else []
    asset_count=len([x for x in assets if x.is_file()])
    if asset_count<case['asset_count'].get('min',0) or asset_count>case['asset_count'].get('max',10**9): errors.append('ASSET_COUNT')
    chunks=[json.loads(x) for x in (bundle/'chunks.jsonl').read_text(encoding='utf-8').splitlines() if x] if (bundle/'chunks.jsonl').exists() else []
    for chunk in chunks:
        if any(value not in idset for value in chunk.get('element_ids',[])) or any(value not in table_ids for value in chunk.get('table_ids',[])): errors.append('REFERENCE_ERROR')
    if any(len(c.get('content',''))>case['max_chunk_chars'] for c in chunks): errors.append('CHUNK_LIMIT')
    if ocr:
        accepted_ids={candidate.get('candidate_id') for candidate in report.get('details',{}).get('ocr_table_candidates',[]) if candidate.get('decision')=='accepted'}
        canonical_ocr=[json.loads(path.read_text(encoding='utf-8')) for path in tables if (json.loads(path.read_text(encoding='utf-8')).get('properties') or {}).get('origin')=='ocr_table_candidate']
        if any((table.get('properties') or {}).get('candidate_id') not in accepted_ids for table in canonical_ocr): errors.append('OCR_TABLE_CONTAINMENT_FAILED')
        if ocr.get('min_accepted_tables',0) and len(canonical_ocr) < ocr['min_accepted_tables']: errors.append('OCR_TABLE_CONTAINMENT_FAILED')
    return errors, report
def _workflow_errors(case, bundle, router_rc, review_rc=0, projection_rc=0):
    errors=[]
    if router_rc: errors.append('WORKFLOW_ROUTER_FAILED')
    if review_rc: errors.append('WORKFLOW_REVIEW_REQUEST_FAILED')
    if projection_rc: errors.append('WORKFLOW_READABLE_PROJECTION_FAILED')
    from validate_bundle import validate_bundle
    try:
        validation=validate_bundle(str(bundle))
    except (OSError, json.JSONDecodeError, TypeError, ValueError):
        errors.append('BUNDLE_VALIDATION_FAILED')
    else:
        if validation.get('status') != 'passed': errors.append('BUNDLE_VALIDATION_FAILED')
    report_path=bundle/'conversion-report.json'; manifest_path=bundle/'manifest.json'; request_path=bundle/'ai-review-request.json'; readable_path=bundle/'document-readable.md'
    report, report_errors=_load_workflow_json(report_path, 'WORKFLOW_REPORT_MISSING', 'WORKFLOW_REPORT_MALFORMED')
    manifest, manifest_errors=_load_workflow_json(manifest_path, 'WORKFLOW_MANIFEST_MISSING', 'WORKFLOW_MANIFEST_MALFORMED')
    request, request_errors=_load_workflow_json(request_path, 'HOST_REVIEW_REQUEST_MISSING', 'WORKFLOW_REQUEST_MALFORMED')
    errors.extend(report_errors + manifest_errors + request_errors)
    if request is not None:
        from ai_review import _schema, fingerprint
        if _schema('ai-review-request.schema.json',request): errors.append('WORKFLOW_REQUEST_SCHEMA_INVALID')
        if manifest is not None and (request.get('source_sha256') != manifest.get('source_sha256') or request.get('canonical_bundle_fingerprint') != fingerprint(bundle)): errors.append('WORKFLOW_REFERENCE_ERROR')
        if not set(case['required_reason_codes']) <= set(request.get('reason_codes',[])): errors.append('WORKFLOW_REASON_CODE_MISSING')
    if report is not None and (report.get('ai_review_request_status') != case['expected_request_status'] or report.get('ai_review_status') != case['expected_ai_review_status'] or report.get('readable_projection_status') != case['expected_readable_projection_status']): errors.append('WORKFLOW_STATUS_MISMATCH')
    if not readable_path.exists(): errors.append('READABLE_PROJECTION_MISSING')
    else:
        text=readable_path.read_text(encoding='utf-8')
        if not text.strip() or any(pattern not in text for pattern in case['required_projection_patterns']): errors.append('READABLE_PROJECTION_SEMANTICS_MISSING')
    return errors

def run_workflow_case(case, root, reruns, keep):
    source=generate(case['fixture'],root/'sources'/case['case_id']); source_sha=hashlib.sha256(Path(source).read_bytes()).hexdigest(); models=[]; fps=[]; errors=[]; bundles=[]
    for run in range(reruns):
        bundle=root/'work'/case['case_id']/str(run); bundle.mkdir(parents=True,exist_ok=True)
        router=subprocess.run([sys.executable,str(ROOT/'scripts/router.py'),str(source),'--output',str(bundle),'--source-url','https://example.test/page/index.html'],text=True,capture_output=True)
        if case.get('workflow') == 'explicit_user_request_review':
            review=subprocess.run([sys.executable,str(ROOT/'scripts/prepare_ai_review.py'),str(bundle),'--force-user-request','--all-eligible-targets'],text=True,capture_output=True)
        else:
            review=subprocess.run([sys.executable,str(ROOT/'scripts/ai_review.py'),str(bundle)],text=True,capture_output=True)
        projection=subprocess.run([sys.executable,str(ROOT/'scripts/render_readable_projection.py'),str(bundle)],text=True,capture_output=True)
        run_errors=_workflow_errors(case,bundle,router.returncode,review.returncode,projection.returncode)
        errors.extend(run_errors)
        # Normalization reads workflow artifacts too. Preserve aggregate output
        # when an artifact contract has already failed instead of reparsing it.
        model={} if any(code in {'WORKFLOW_REPORT_MISSING','WORKFLOW_REPORT_MALFORMED','WORKFLOW_MANIFEST_MISSING','WORKFLOW_MANIFEST_MALFORMED','HOST_REVIEW_REQUEST_MISSING','WORKFLOW_REQUEST_MALFORMED','WORKFLOW_ARTIFACT_UNREADABLE'} for code in run_errors) else normalize_bundle(bundle)
        models.append(model); fps.append(fingerprint(model)); bundles.append(bundle)
    if any(not semantically_equal(models[0],model) for model in models[1:]): errors.append('WORKFLOW_NONDETERMINISTIC_RERUN'); write_diff(models[0],models[1],root/'diffs',case['case_id'])
    result={'case_id':case['case_id'],'format':'workflow','workflow':case['workflow'],'status':'passed' if not errors else 'failed','reason_codes':sorted(set(errors)),'fixture_recipe_id':case['fixture'],'fixture_description':'checked-in HTML fixture routed through the production Phase 5 workflow','generated_source_sha256':source_sha,'fingerprints':fps,'normalized_snapshot':models[0]}
    if result['status']!='passed' or keep: result['bundles']=[str(p) for p in bundles]
    return result

def run_case(case, root, reruns, keep):
    if case['profile']=='optional-pandoc' and not shutil.which('pandoc'): return {'case_id':case['case_id'],'format':case['format'],'status':'skipped','skip_reason':'pandoc binary unavailable'}
    source=generate(case['fixture'],root/'sources'/case['case_id']); source_sha=hashlib.sha256(Path(source).read_bytes()).hexdigest(); models=[]; fps=[]; errors=[]; bundles=[]
    for run in range(reruns):
        bundle=root/'work'/case['case_id']/str(run); bundle.mkdir(parents=True,exist_ok=True)
        args=[sys.executable,str(ROOT/'scripts/router.py'),str(source),'--output',str(bundle)]
        if case['format']=='html': args += ['--source-url','https://example.test/page/index.html']
        process=subprocess.run(args,text=True,capture_output=True); current,report=assert_contract(case,bundle,process.returncode); errors += current; models.append(normalize_bundle(bundle));fps.append(fingerprint(models[-1]));bundles.append(bundle)
    rerun_mismatch=any(not semantically_equal(models[0],model) for model in models[1:])
    if rerun_mismatch: errors.append('CROSS_FORMAT_NONDETERMINISTIC_RERUN');write_diff(models[0],models[1],root/'diffs',case['case_id'])
    status='passed' if not errors else 'failed'; result={'case_id':case['case_id'],'format':case['format'],'status':status,'reason_codes':sorted(set(errors)),'fixture_recipe_id':case['fixture'],'fixture_description':'programmatically generated regression source','generator_dependency':'repository test dependencies','generated_source_sha256':source_sha,'fingerprints':fps,'normalized_snapshot':models[0]}
    if status!='passed' or keep: result['bundles']=[str(p) for p in bundles]
    return result
def main(args):
    if not (ROOT / "tests").is_dir() or not MANIFEST.is_file():
        print(
            "This regression runner requires the complete source repository "
            "and cannot run from the Release Package or Agent Skill archive.",
            file=sys.stderr,
        )
        return 2

    if args.update_baseline and not args.confirm_baseline_update: raise SystemExit('--update-baseline requires --confirm-baseline-update')
    all_format_cases=load_cases(); all_workflow_cases=load_workflow_cases()
    duplicate_case_ids={c['case_id'] for c in all_format_cases} & {c['case_id'] for c in all_workflow_cases}
    if duplicate_case_ids: raise AssertionError('duplicate case_id across format and workflow cases: '+', '.join(sorted(duplicate_case_ids)))
    format_cases=[c for c in all_format_cases if c['profile']==args.profile and (not args.case or c['case_id']==args.case) and (not args.format or c['format']==args.format)]
    workflow_cases=[c for c in all_workflow_cases if c['profile']==args.profile and (not args.case or c['case_id']==args.case) and not args.format]
    cases=format_cases+workflow_cases
    output=Path(args.output);output.mkdir(parents=True,exist_ok=True);env=environment_manifest('pandoc-enabled' if args.profile=='optional-pandoc' else 'core-no-pandoc');(output/'environment-manifest.json').write_text(json.dumps(env,indent=2)+'\n')
    results=[]
    for case in cases:
        print(f"[cross-format] starting {case['case_id']}", file=sys.stderr, flush=True)
        result=run_case(case,output,args.reruns,args.keep_bundles) if 'format' in case else run_workflow_case(case,output,args.reruns,args.keep_bundles)
        results.append(result)
        print(f"[cross-format] finished {case['case_id']}: {result['status']}", file=sys.stderr, flush=True)
    baseline_dir=Path(args.baseline_dir); baseline_path=baseline_dir/'fingerprints.json'; baseline=json.loads(baseline_path.read_text(encoding='utf-8')) if baseline_path.exists() else {}
    baseline_mismatches=[]
    for result in results:
        current=result.get('fingerprints',[{}])[0].get('bundle_fingerprint')
        expected=baseline.get(result['case_id'])
        snapshot=baseline_dir/(result['case_id']+'.normalized.json')
        baseline_changed=expected and expected != current
        # A fingerprint can differ solely because an OCR confidence moved within
        # the declared tolerance; snapshots make that distinction auditable.
        if baseline_changed and snapshot.exists():
            baseline_changed=not semantically_equal(json.loads(snapshot.read_text(encoding='utf-8')),result['normalized_snapshot'])
        if result['status']=='passed' and baseline_changed:
            result['reason_codes'].append('BASELINE_FINGERPRINT_MISMATCH'); baseline_mismatches.append(result['case_id'])
            if not args.update_baseline: result['status']='failed'
            if snapshot.exists(): write_diff(json.loads(snapshot.read_text(encoding='utf-8')),result['normalized_snapshot'],output/'diffs',result['case_id']+'-baseline')
    if args.update_baseline:
        if any(r['status']!='passed' for r in results): raise SystemExit('baseline update refused: all selected cases must pass contract validation')
        updated={r['case_id']:r['fingerprints'][0]['bundle_fingerprint'] for r in results}
        merged=dict(baseline); merged.update(updated)
        baseline_path.parent.mkdir(parents=True,exist_ok=True)
        report={'old':baseline,'new':merged,'changed':[key for key in updated if baseline.get(key)!=updated[key]],'preserved':[key for key in baseline if key not in updated]}
        (output/'baseline-update-report.json').write_text(json.dumps(report,indent=2)+'\n'); baseline_path.write_text(json.dumps(merged,indent=2)+'\n')
        for result in results: (baseline_dir/(result['case_id']+'.normalized.json')).write_text(json.dumps(result['normalized_snapshot'],ensure_ascii=False,indent=2)+'\n')
    counts=Counter(r['status'] for r in results); formats=defaultdict(lambda:Counter())
    for r in results: formats[r['format']][r['status']]+=1
    def has_rerun_mismatch(result): return bool({'CROSS_FORMAT_NONDETERMINISTIC_RERUN','WORKFLOW_NONDETERMINISTIC_RERUN'} & set(result.get('reason_codes',[])))
    rerun_mismatches=sum(has_rerun_mismatch(r) for r in results)
    summary={'schema_version':'1.0','task':f'REPRODUCIBLE-CROSS-FORMAT-REGRESSION-{SKILL_VERSION.upper()}','profile':args.profile,'environment_profile':env['profile'],'format_case_count':len(format_cases),'workflow_case_count':len(workflow_cases),'case_count':len(results),'passed':counts['passed'],'failed':counts['failed'],'skipped':counts['skipped'],'formats':{k:dict(v) for k,v in formats.items()},'bundle_validation_failures':sum(any('BUNDLE_VALIDATION' in x for x in r.get('reason_codes',[])) for r in results),'unexpected_warning_cases':sum('UNEXPECTED_WARNING' in r.get('reason_codes',[]) for r in results),'chunk_limit_violations':sum('CHUNK_LIMIT' in r.get('reason_codes',[]) for r in results),'reference_errors':sum(any(x in r.get('reason_codes',[]) for x in ('REFERENCE_ERROR','ELEMENT_ID_INVALID','LOCATOR_MISSING','TABLE_INDEX_MALFORMED')) for r in results),'rerun_mismatches':rerun_mismatches,'baseline_mismatches':len(baseline_mismatches),'workflow_failures':sum(r.get('format')=='workflow' and r['status']=='failed' for r in results),'normalized_determinism_status':'passed' if not rerun_mismatches else 'failed','validation_status':'passed' if not counts['failed'] else 'failed','cases':results}
    public_results=[{k:v for k,v in result.items() if k!='normalized_snapshot'} for result in results]
    summary['cases']=public_results
    (output/'cross-format-regression-cases.jsonl').write_text(''.join(json.dumps(r,ensure_ascii=False)+'\n' for r in public_results));(output/'artifact-fingerprints.json').write_text(json.dumps({r['case_id']:r.get('fingerprints',[]) for r in results},indent=2)+'\n');(output/'cross-format-regression-summary.json').write_text(json.dumps(summary,indent=2)+'\n');print(json.dumps(summary,indent=2));return 1 if counts['failed'] else 0
if __name__=='__main__':
 p=argparse.ArgumentParser();p.add_argument('--output',required=True);p.add_argument('--profile',default='core',choices=['core','optional-pandoc']);p.add_argument('--case');p.add_argument('--format');p.add_argument('--reruns',type=int,default=2);p.add_argument('--baseline-dir',default=str(ROOT/'tests/cross_format/baselines'));p.add_argument('--update-baseline',action='store_true');p.add_argument('--confirm-baseline-update',action='store_true');p.add_argument('--keep-bundles',action='store_true');a=p.parse_args();sys.exit(main(a))
