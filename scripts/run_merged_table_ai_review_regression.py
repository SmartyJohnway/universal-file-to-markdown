#!/usr/bin/env python3
"""Run true cross-format merged-table AI Review trigger regression.

This focused source-repository regression generates merged and plain tables for
DOCX, XLSX, PPTX, and HTML, routes them through the production converter, then
invokes prepare_ai_review.py in both automatic and explicit target-table modes.
Positive cases must generate table-targeted AI Review requests with correct
truthful reason codes. Plain-table controls must not report geometry reasons.
Canonical bundle files must remain byte-identical before and after request preparation.
"""

from __future__ import annotations

import argparse
import hashlib
import json
import subprocess
import sys
import tempfile
from pathlib import Path
from typing import Any, Callable

ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "scripts"
CANONICAL_NAMES = ("document.json", "chunks.jsonl", "manifest.json")
MERGED_REASON = "MERGED_TABLE_GEOMETRY_PRESENT"
HTML_COMPLEX_REASON = "HTML_MERGED_TABLE_COMPLEX"
EXPLICIT_REASON = "EXPLICIT_USER_REQUEST"


def _sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()


def _canonical_hashes(bundle: Path) -> dict[str, str]:
    paths = [bundle / name for name in CANONICAL_NAMES]
    tables_dir = bundle / "tables"
    if tables_dir.is_dir():
        paths.extend(sorted(p for p in tables_dir.glob("*.json")))
    return {str(path.relative_to(bundle)): _sha256(path) for path in paths if path.is_file()}


def _load_json(path: Path) -> dict[str, Any]:
    value = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(value, dict):
        raise AssertionError(f"Expected JSON object: {path}")
    return value


def _table_ids(bundle: Path) -> set[str]:
    result: set[str] = set()
    tables_dir = bundle / "tables"
    if not tables_dir.is_dir():
        return result
    for path in tables_dir.glob("*.json"):
        if path.name == "index.json":
            continue
        value = _load_json(path)
        table_id = value.get("id")
        if isinstance(table_id, str) and table_id:
            result.add(table_id)
    return result


def _write_docx(path: Path, merged: bool) -> None:
    from docx import Document

    document = Document()
    document.add_paragraph("Merged table trigger regression")
    table = document.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Header A"
    table.cell(0, 1).text = "Header B"
    table.cell(1, 0).text = "Value A"
    table.cell(1, 1).text = "Value B"
    if merged:
        table.cell(0, 0).merge(table.cell(0, 1))
        table.cell(0, 0).text = "Merged Header"
    document.save(path)


def _write_xlsx(path: Path, merged: bool) -> None:
    import openpyxl

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet["A1"] = "Header A"
    sheet["B1"] = "Header B"
    sheet["A2"] = "Value A"
    sheet["B2"] = "Value B"
    if merged:
        sheet.merge_cells("A1:B1")
        sheet["A1"] = "Merged Header"
    workbook.save(path)


def _write_pptx(path: Path, merged: bool) -> None:
    from pptx import Presentation
    from pptx.util import Inches

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(5), Inches(2))
    table = shape.table
    table.cell(0, 0).text = "Header A"
    table.cell(0, 1).text = "Header B"
    table.cell(1, 0).text = "Value A"
    table.cell(1, 1).text = "Value B"
    if merged:
        table.cell(0, 0).merge(table.cell(0, 1))
        table.cell(0, 0).text = "Merged Header"
    presentation.save(path)


def _write_html(path: Path, merged: bool) -> None:
    header = '<th colspan="2">Merged Header</th>' if merged else "<th>Header A</th><th>Header B</th>"
    path.write_text(
        "<!doctype html><html><head><title>Regression</title></head><body>"
        "<main><h1>Merged table trigger regression</h1><table><tr>"
        + header
        + "</tr><tr><td>Value A</td><td>Value B</td></tr></table></main>"
        "</body></html>",
        encoding="utf-8",
    )


WRITERS: dict[str, tuple[str, Callable[[Path, bool], None]]] = {
    "docx": ("docx", _write_docx),
    "xlsx": ("xlsx", _write_xlsx),
    "pptx": ("pptx", _write_pptx),
    "html": ("html", _write_html),
}


def _run(command: list[str]) -> subprocess.CompletedProcess[str]:
    return subprocess.run(command, cwd=ROOT, text=True, capture_output=True)


def _case(format_name: str, merged: bool, mode: str, root: Path) -> dict[str, Any]:
    suffix, writer = WRITERS[format_name]
    case_id = f"{format_name}-{'merged' if merged else 'plain'}-{mode}-ai-review-trigger"
    source_dir = root / "sources" / case_id
    bundle = root / "bundles" / case_id
    source_dir.mkdir(parents=True, exist_ok=True)
    bundle.mkdir(parents=True, exist_ok=True)
    source = source_dir / f"source.{suffix}"
    writer(source, merged)

    router_command = [sys.executable, str(SCRIPTS / "router.py"), str(source), "--output", str(bundle)]
    if format_name == "html":
        router_command.extend(["--source-url", "https://example.test/merged/index.html"])
    router = _run(router_command)

    errors: list[str] = []
    if router.returncode != 0:
        errors.append("ROUTER_FAILED")
    required = (bundle / "conversion-report.json", bundle / "document.json", bundle / "manifest.json")
    if any(not path.is_file() for path in required):
        errors.append("BUNDLE_MISSING")
        return {
            "case_id": case_id,
            "format": format_name,
            "merged": merged,
            "mode": mode,
            "status": "failed",
            "reason_codes": errors,
            "router_stderr": router.stderr,
        }

    canonical_table_ids = _table_ids(bundle)
    before = _canonical_hashes(bundle)

    prepare_cmd = [sys.executable, str(SCRIPTS / "prepare_ai_review.py"), str(bundle)]
    if mode == "explicit":
        if not canonical_table_ids:
            errors.append("CANONICAL_TABLE_MISSING_FOR_EXPLICIT_TEST")
            return {
                "case_id": case_id,
                "format": format_name,
                "merged": merged,
                "mode": mode,
                "status": "failed",
                "reason_codes": errors,
            }
        target_id = sorted(canonical_table_ids)[0]
        prepare_cmd.extend(["--force-user-request", "--target-table", target_id])

    prepare = _run(prepare_cmd)
    after = _canonical_hashes(bundle)
    if prepare.returncode != 0:
        errors.append("PREPARE_AI_REVIEW_FAILED")
    if before != after:
        errors.append("CANONICAL_MUTATION")

    report = _load_json(bundle / "conversion-report.json")
    assessment = report.get("quality_risk_assessment") or {}
    assessment_reasons = set(assessment.get("reason_codes") or [])
    request_path = bundle / "ai-review-request.json"
    request = _load_json(request_path) if request_path.is_file() else None
    request_reasons = set(request.get("reason_codes") or []) if request else set()
    targets = request.get("targets") or [] if request else []

    # Verification invariants
    if mode == "automatic":
        if merged:
            if not report.get("ai_review_recommended"):
                errors.append("AI_REVIEW_NOT_RECOMMENDED")
            if MERGED_REASON not in assessment_reasons:
                errors.append("ASSESSMENT_REASON_MISSING")
            if request is None:
                errors.append("AI_REVIEW_REQUEST_MISSING")
            else:
                if MERGED_REASON not in request_reasons:
                    errors.append("REQUEST_REASON_MISSING")
                matching_targets = [
                    t
                    for t in targets
                    if t.get("target_type") == "table" and t.get("target_id") in canonical_table_ids
                ]
                if not matching_targets:
                    errors.append("CANONICAL_TABLE_TARGET_MISSING")
                if any(MERGED_REASON not in set(t.get("reason_codes") or []) for t in matching_targets):
                    errors.append("TARGET_MERGED_REASON_MISSING")
        else:
            if MERGED_REASON in assessment_reasons:
                errors.append("PLAIN_TABLE_FALSE_POSITIVE_ASSESSMENT")
            if MERGED_REASON in request_reasons:
                errors.append("PLAIN_TABLE_FALSE_POSITIVE_REQUEST")
            if any(MERGED_REASON in set(t.get("reason_codes") or []) for t in targets):
                errors.append("PLAIN_TABLE_FALSE_POSITIVE_TARGET")
            if HTML_COMPLEX_REASON in assessment_reasons or HTML_COMPLEX_REASON in request_reasons:
                errors.append("PLAIN_TABLE_FALSE_HTML_COMPLEX_REASON")

    elif mode == "explicit":
        if request is None:
            errors.append("EXPLICIT_REQUEST_NOT_GENERATED")
        else:
            if EXPLICIT_REASON not in request_reasons:
                errors.append("EXPLICIT_REASON_MISSING_FROM_REQUEST")
            matching_targets = [
                t
                for t in targets
                if t.get("target_type") == "table" and t.get("target_id") in canonical_table_ids
            ]
            if not matching_targets:
                errors.append("EXPLICIT_TARGET_MISSING")
            else:
                for target in matching_targets:
                    t_reasons = set(target.get("reason_codes") or [])
                    if EXPLICIT_REASON not in t_reasons:
                        errors.append("EXPLICIT_REASON_MISSING_FROM_TARGET")
                    if merged:
                        if MERGED_REASON not in t_reasons:
                            errors.append("MERGED_REASON_MISSING_FROM_EXPLICIT_TARGET")
                        if format_name == "html" and HTML_COMPLEX_REASON not in t_reasons:
                            errors.append("HTML_COMPLEX_REASON_MISSING_FROM_EXPLICIT_TARGET")
                        if format_name != "html" and HTML_COMPLEX_REASON in t_reasons:
                            errors.append("HTML_COMPLEX_REASON_FALSE_POSITIVE_IN_EXPLICIT_TARGET")
                    else:
                        if MERGED_REASON in t_reasons:
                            errors.append("PLAIN_EXPLICIT_TARGET_FALSE_MERGED_REASON")
                        if HTML_COMPLEX_REASON in t_reasons:
                            errors.append("PLAIN_EXPLICIT_TARGET_FALSE_HTML_COMPLEX_REASON")

    return {
        "case_id": case_id,
        "format": format_name,
        "merged": merged,
        "mode": mode,
        "status": "passed" if not errors else "failed",
        "reason_codes": sorted(set(errors)),
        "router_returncode": router.returncode,
        "prepare_returncode": prepare.returncode,
        "assessment_reason_codes": sorted(assessment_reasons),
        "request_reason_codes": sorted(request_reasons),
        "canonical_table_ids": sorted(canonical_table_ids),
        "request_target_ids": sorted(
            target.get("target_id") for target in targets if isinstance(target.get("target_id"), str)
        ),
        "canonical_hashes_preserved": before == after,
        "router_stderr": router.stderr,
        "prepare_stderr": prepare.stderr,
    }


def run(output: Path) -> dict[str, Any]:
    output.mkdir(parents=True, exist_ok=True)
    results = [
        _case(format_name, merged, mode, output)
        for format_name in WRITERS
        for merged in (True, False)
        for mode in ("automatic", "explicit")
    ]
    failed = [result for result in results if result["status"] != "passed"]
    summary = {
        "suite": "merged-table-ai-review-trigger-e2e",
        "case_count": len(results),
        "positive_case_count": sum(1 for result in results if result["merged"] or result["mode"] == "explicit"),
        "negative_case_count": sum(1 for result in results if not result["merged"] and result["mode"] == "automatic"),
        "passed": len(results) - len(failed),
        "failed": len(failed),
        "validation_status": "passed" if not failed else "failed",
        "results": results,
    }
    (output / "merged-table-ai-review-regression-report.json").write_text(
        json.dumps(summary, ensure_ascii=False, indent=2) + "\n", encoding="utf-8"
    )
    return summary


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--output", help="Evidence output directory")
    args = parser.parse_args()
    if not (ROOT / "tests").is_dir():
        print(
            "This regression runner requires the complete source repository and cannot run "
            "from the Release Package or Agent Skill archive.",
            file=sys.stderr,
        )
        return 2

    if args.output:
        output = Path(args.output).resolve()
        summary = run(output)
    else:
        with tempfile.TemporaryDirectory(prefix="ufm-merged-trigger-") as tmp_dir:
            summary = run(Path(tmp_dir))
    print(json.dumps(summary, ensure_ascii=False, indent=2))
    return 0 if summary["validation_status"] == "passed" else 1


if __name__ == "__main__":
    raise SystemExit(main())
