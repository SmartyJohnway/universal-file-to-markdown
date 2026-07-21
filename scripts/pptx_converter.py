"""
pptx_converter.py (v1.5.1)
Custom PowerPoint -> Markdown converter, replacing the v1.1-v1.4
MarkItDown fallback. python-pptx gives a clean, first-class API for the
things this skill cares about most:
  - table merges: cell.is_merge_origin / is_spanned / span_width /
    span_height - no raw-XML digging needed, unlike DOCX/XLSX
  - chart title/categories/series via shape.chart
  - speaker notes via slide.notes_slide

Reading order: shapes are sorted by (top, left) position, which is a
reasonable approximation of visual reading order for typical slide
layouts (title/body/footer stacked top-to-bottom). This is a heuristic,
not a guarantee - a slide with an intentionally unusual spatial layout
(e.g. two independent side-by-side columns meant to be read as separate
flows) can still come out in an order that doesn't match the presenter's
intent. See engine_notes.md.

v1.5.1 fixes (see references/engine_notes.md for the bugs these replace):
  - level-0 bullets are no longer treated as plain paragraphs. Bullet-ness
    is read from the paragraph's own <a:pPr> (a:buChar / a:buAutoNum /
    a:buNone), not guessed from indentation level.
  - SmartArt diagrams and embedded OLE objects are no longer a silent
    drop: the pptx zip container's relationship parts are scanned so
    every input surfaces SMARTART_NOT_EXTRACTED / EMBEDDED_OLE_NOT_EXTRACTED
    in the conversion report when present, instead of only being
    documented as an abstract "known limitation".
  - group shapes now collect ALL nested tables (previously only the
    first, and only in a comment - the actual code discarded every
    group-nested table's standalone asset).
  - merged-cell tables now emit a `cells` list (row/col/rowspan/colspan)
    alongside the flattened `rows` grid, so tables/*.html keeps the span
    geometry instead of being silently flattened by table_export.py.

Still deliberately out of scope for this version (see engine_notes.md):
  - SmartArt/OLE CONTENT extraction (text inside the diagram/object) -
    only detection + disclosure is in scope; the content itself is still
    not extracted.
  - Precise text run-level bold/italic within slide text boxes (slide text
    is usually short titles/bullets where this matters less than in a
    Word document body; can be added later following the docx pattern)
"""

import os
import re
import zipfile
from xml.etree import ElementTree as ET

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from common_utils import extract_ooxml_core_metadata

_REL_NS = "http://schemas.openxmlformats.org/package/2006/relationships"
_OLE_REL_SUFFIX = "/relationships/oleObject"
_DIAGRAM_REL_SUFFIX = "/relationships/diagramData"


def convert_pptx(path: str, assets_dir: str = None) -> dict:
    prs = Presentation(path)

    sections = []
    elements = []
    tables_out = []
    table_counter = 0
    image_count = 0
    chart_count = 0
    notes_count = 0

    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        slide_id = f"slide-{slide_num:04d}"
        blocks = [f"<!-- slide: {slide_num} -->"]
        slide_children = []

        top_shapes = sorted(slide.shapes, key=_shape_sort_key)
        for shape in top_shapes:
            rendered = _render_shape(shape, assets_dir, slide_num)
            if rendered["markdown"]:
                blocks.append(rendered["markdown"])
            image_count += rendered["image_count"]
            chart_count += rendered["chart_count"]
            assigned_table_ids = []
            table_locators = [item.get("source_locator", {}) for item in rendered["items"]
                              if item.get("type") == "table"]
            for table_offset, tbl in enumerate(rendered["tables"]):
                table_counter += 1
                tbl["id"] = f"table-{table_counter:04d}"
                tbl["context"] = f"pptx_slide_{slide_num}"
                tbl["source_locator"] = {"slide": slide_num,
                                         **(table_locators[table_offset]
                                            if table_offset < len(table_locators) else {})}
                tbl["engine"] = "python-pptx_custom"
                assigned_table_ids.append(tbl["id"])
                tables_out.append(tbl)
            table_id_iter = iter(assigned_table_ids)
            for item in rendered["items"]:
                if item.get("type") == "table":
                    item["table_id"] = next(table_id_iter, None)
                slide_children.append(item)

        if slide.has_notes_slide:
            notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
            if notes_text:
                notes_count += 1
                blocks.append(f"<!-- speaker notes -->\n{notes_text}")
                slide_children.append({
                    "type": "speaker_note", "content": notes_text,
                    "token": f"notes-{slide_num}",
                    "source_locator": {"slide": slide_num, "part": "speaker_notes"},
                })

        slide_md = "\n\n".join(b for b in blocks if b)
        sections.append(slide_md)
        elements.append({
            "id": slide_id,
            "type": "slide",
            "slide": slide_num,
            "content": f"<!-- slide: {slide_num} -->",
            "engine": "python-pptx_custom",
            "confidence": None,
            "source_locator": {"slide": slide_num},
        })
        token_to_id = {item.get("token"): f"{slide_id}-{item.get('token')}"
                       for item in slide_children if item.get("token")}
        for child_idx, raw_item in enumerate(slide_children, start=1):
            item = dict(raw_item)
            item_type = item.pop("type")
            token = item.pop("token", None)
            parent_token = item.pop("parent_token", None)
            table_id = item.pop("table_id", None)
            locator = {"slide": slide_num}
            locator.update(item.pop("source_locator", {}) or {})
            element = {
                "id": token_to_id.get(token, f"{slide_id}-{child_idx:03d}"),
                "parent_id": token_to_id.get(parent_token, slide_id),
                "type": item_type,
                "content": item.pop("content", ""),
                "engine": "python-pptx_custom",
                "confidence": None,
                "source_locator": locator,
            }
            element.update(item)
            if table_id:
                element["table_id"] = table_id
            elements.append(element)

    special_parts = _scan_special_parts(path)

    metadata = extract_ooxml_core_metadata(path)
    report = {
        "status": "passed",
        "engine": "python-pptx_custom",
        "slide_count": len(prs.slides._sldIdLst),
        "tables_found": table_counter,
        "images_found": image_count,
        "charts_found": chart_count,
        "speaker_notes_found": notes_count,
        "media_extracted": image_count,
        "smartart_parts_found": special_parts["smartart_parts_found"],
        "ole_objects_found": special_parts["ole_objects_found"],
        "smartart_occurrences": special_parts["smartart_occurrences"],
        "ole_occurrences": special_parts["ole_occurrences"],
        "metadata": metadata,
    }
    if special_parts["smartart_parts_found"] or special_parts["ole_objects_found"]:
        report["status"] = "passed_with_warnings"

    return {"markdown": "\n\n".join(sections), "report": report,
            "elements": elements, "tables": tables_out}


def _scan_special_parts(path: str) -> dict:
    """Scan the pptx zip container's relationship parts directly, rather
    than relying on python-pptx (which has no first-class SmartArt/OLE
    API), so SmartArt diagrams and embedded OLE objects are DETECTED and
    DISCLOSED in the report even though their content still isn't
    extracted. Previously these were only mentioned in SKILL.md as an
    abstract "known limitation" with nothing in conversion-report.json
    reflecting whether a given input actually contained any - i.e. every
    real file with SmartArt/OLE looked identical to one without any."""
    smartart_count = 0
    ole_count = 0
    smartart_occurrences = []
    ole_occurrences = []
    try:
        with zipfile.ZipFile(path) as z:
            names = z.namelist()
            # corroborating evidence: dedicated parts under ppt/diagrams/
            smartart_parts = {n for n in names if n.startswith("ppt/diagrams/") and n.endswith(".xml")}

            rels_files = [n for n in names if n.startswith("ppt/slides/_rels/") and n.endswith(".rels")]
            for rels_name in rels_files:
                try:
                    tree = ET.fromstring(z.read(rels_name))
                except ET.ParseError:
                    continue
                for rel in tree.findall(f"{{{_REL_NS}}}Relationship"):
                    rtype = rel.get("Type", "")
                    match = re.search(r"/slide(\d+)\.xml\.rels$", rels_name)
                    slide_number = int(match.group(1)) if match else None
                    occurrence = {"slide": slide_number, "relationship_id": rel.get("Id"),
                                  "target": rel.get("Target"), "relationship_part": rels_name}
                    if rtype.endswith(_DIAGRAM_REL_SUFFIX):
                        smartart_count += 1
                        smartart_occurrences.append(occurrence)
                    elif rtype.endswith(_OLE_REL_SUFFIX):
                        ole_count += 1
                        ole_occurrences.append(occurrence)

            if smartart_count == 0 and smartart_parts:
                # diagrams/ parts exist but no slide rel pointed at them
                # (unusual, but don't under-report if the folder is there)
                smartart_count = len(smartart_parts)
                smartart_occurrences = [{"slide": None, "relationship_id": None,
                                         "target": part, "relationship_part": None}
                                        for part in sorted(smartart_parts)]
    except (zipfile.BadZipFile, FileNotFoundError, OSError):
        pass

    return {"smartart_parts_found": smartart_count, "ole_objects_found": ole_count,
            "smartart_occurrences": smartart_occurrences,
            "ole_occurrences": ole_occurrences}


def _shape_sort_key(shape):
    top = shape.top if shape.top is not None else 0
    left = shape.left if shape.left is not None else 0
    return (top, left)


def _render_shape(shape, assets_dir, slide_num, token_prefix="") -> dict:
    """Returns {"markdown": str, "image_count": int, "chart_count": int,
    "tables": [table_dict, ...]}. `tables` is a list (not a single
    optional dict) specifically so group shapes can propagate every
    nested table's standalone asset, not just the first."""
    empty = {"markdown": "", "image_count": 0, "chart_count": 0, "tables": [], "items": []}
    shape_token = f"{token_prefix}shape-{shape.shape_id}"

    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        pieces = []
        img_c = chart_c = 0
        tables = []
        items = []
        group_token = f"{token_prefix}group-{shape.shape_id}"
        group_item = {"type": "group", "content": "", "token": group_token,
                      "source_locator": _shape_locator(shape)}
        for sub in sorted(shape.shapes, key=_shape_sort_key):
            sub_result = _render_shape(sub, assets_dir, slide_num,
                                       token_prefix=f"{group_token}-")
            if sub_result["markdown"]:
                pieces.append(sub_result["markdown"])
            img_c += sub_result["image_count"]
            chart_c += sub_result["chart_count"]
            tables.extend(sub_result["tables"])
            for item in sub_result["items"]:
                if not item.get("parent_token"):
                    item["parent_token"] = group_token
                items.append(item)
        return {"markdown": "\n\n".join(pieces), "image_count": img_c,
                "chart_count": chart_c, "tables": tables,
                "items": [group_item] + items}

    if getattr(shape, "has_chart", False):
        markdown = _render_chart(shape.chart)
        return {"markdown": markdown, "image_count": 0,
                "chart_count": 1, "tables": [],
                "items": [{"type": "chart", "content": markdown,
                           "token": shape_token,
                           "source_locator": {**_shape_locator(shape),
                                              "part": str(shape.chart.part.partname)}}]}

    if shape.has_table:
        md, grid, cells = _render_table(shape.table)
        table_dict = {"rows": grid, "context": "pptx_table"}
        if cells is not None:
            table_dict["cells"] = cells
        return {"markdown": md, "image_count": 0, "chart_count": 0,
                "tables": [table_dict],
                "items": [{"type": "table", "content": md,
                           "token": shape_token,
                           "source_locator": _shape_locator(shape)}]}

    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        asset, relationship_id, part = _extract_picture_asset(shape, assets_dir, slide_num)
        if asset:
            alt = (shape.name or "").strip()
            markdown = f"![{alt}]({asset})"
            return {"markdown": markdown, "image_count": 1,
                    "chart_count": 0, "tables": [],
                    "items": [{"type": "image", "content": markdown,
                               "token": shape_token, "asset": asset,
                               "source_locator": {**_shape_locator(shape),
                                                  "relationship_id": relationship_id,
                                                  "part": part}}]}
        return dict(empty)

    if shape.has_text_frame:
        is_title = _is_title_placeholder(shape)
        text = _render_text_frame(shape, is_title=is_title)
        if not text:
            return dict(empty)
        markdown = f"# {text}" if is_title else text
        semantic_type = "title" if is_title else (
            "list" if any(re.match(r"^\s*(?:-|\d+\.)\s+", line)
                          for line in markdown.splitlines()) else "paragraph")
        return {"markdown": markdown, "image_count": 0,
                "chart_count": 0, "tables": [],
                "items": [{"type": semantic_type,
                           "content": markdown, "token": shape_token,
                           "source_locator": _shape_locator(shape)}]}

    return dict(empty)


def _shape_locator(shape) -> dict:
    return {
        "shape_id": getattr(shape, "shape_id", None),
        "shape_name": getattr(shape, "name", None),
        "bbox": {
            "left": getattr(shape, "left", None), "top": getattr(shape, "top", None),
            "width": getattr(shape, "width", None), "height": getattr(shape, "height", None),
        },
    }


def _extract_picture_asset(shape, assets_dir, slide_num):
    if not assets_dir:
        return None, None, None
    try:
        image = shape.image
        ext = image.ext or "bin"
        filename = f"slide-{slide_num:04d}-shape-{shape.shape_id:04d}-image.{ext}"
        os.makedirs(assets_dir, exist_ok=True)
        with open(os.path.join(assets_dir, filename), "wb") as f:
            f.write(image.blob)
        blip = shape._element.blipFill.blip
        relationship_id = blip.rEmbed
        rel = shape.part.rels[relationship_id]
        part = str(rel.target_part.partname)
        return filename, relationship_id, part
    except Exception:
        return None, None, None


def _is_title_placeholder(shape) -> bool:
    if not shape.is_placeholder:
        return False
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER
        return shape.placeholder_format.type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
    except Exception:
        return False


def _bullet_from_properties(pPr, level: int):
    if pPr is None:
        return None
    if pPr.find(qn("a:buNone")) is not None:
        return {"is_bullet": False, "ordered": False, "level": level, "char": None}
    if pPr.find(qn("a:buAutoNum")) is not None:
        return {"is_bullet": True, "ordered": True, "level": level, "char": None}
    bu_char = pPr.find(qn("a:buChar"))
    if bu_char is not None:
        return {"is_bullet": True, "ordered": False, "level": level,
                "char": bu_char.get("char") or "\u2022"}
    return None


def _is_body_like_placeholder(shape) -> bool:
    if not getattr(shape, "is_placeholder", False):
        return False
    try:
        from pptx.enum.shapes import PP_PLACEHOLDER
        allowed = {PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT}
        for name in ("TEXT", "VERTICAL_BODY", "VERTICAL_OBJECT"):
            value = getattr(PP_PLACEHOLDER, name, None)
            if value is not None:
                allowed.add(value)
        return shape.placeholder_format.type in allowed
    except Exception:
        return False


def _level_properties(container, level: int):
    if container is None:
        return None
    return container.find(qn(f"a:lvl{min(max(level, 0), 8) + 1}pPr"))


def _resolve_inherited_bullet(shape, level: int):
    """Resolve the common DrawingML inheritance chain conservatively:
    shape list style -> matching layout placeholder -> master body style.
    Return None when the chain carries no explicit bullet declaration."""
    tx_body = getattr(shape._element, "txBody", None)
    result = _bullet_from_properties(
        _level_properties(getattr(tx_body, "lstStyle", None), level), level)
    if result is not None:
        return result

    try:
        idx = shape.placeholder_format.idx
        layout = shape.part.slide_layout
        layout_shape = next((ph for ph in layout.placeholders
                             if ph.placeholder_format.idx == idx), None)
        if layout_shape is not None:
            layout_tx_body = getattr(layout_shape._element, "txBody", None)
            result = _bullet_from_properties(
                _level_properties(getattr(layout_tx_body, "lstStyle", None), level), level)
            if result is not None:
                return result

        master = layout.slide_master
        tx_styles = master._element.find(qn("p:txStyles"))
        body_style = tx_styles.find(qn("p:bodyStyle")) if tx_styles is not None else None
        result = _bullet_from_properties(_level_properties(body_style, level), level)
        if result is not None:
            return result
    except Exception:
        pass
    return None


def _get_bullet_info(paragraph, shape) -> dict:
    """Read bullet-ness from the paragraph's own <a:pPr> element instead
    of inferring it from indentation `level`. The v1.5 code treated
    `level == 0` as "not a bullet", but level 0 just means "first
    indentation level" - PowerPoint's ordinary top-level bulleted text is
    level 0. That collapsed the most common case (a flat bulleted list)
    into plain paragraphs, losing list semantics entirely.

    This reads the paragraph's explicit bullet markup:
      - <a:buNone/>            -> not a bullet
      - <a:buAutoNum .../>      -> numbered list item
      - <a:buChar char="..."/>  -> bulleted list item (using that char)
      - none of the above       -> no explicit override on this paragraph;
                                   body-placeholder text defaults to
                                   bulleted in PowerPoint's own list
                                   style, so this still renders as a
                                   bullet. Full inheritance resolution
                                   from the slide layout/master bullet
                                   definitions is not attempted (out of
                                   scope - see engine_notes.md); a
                                   paragraph that explicitly overrides
                                   bullets via buNone/buChar/buAutoNum is
                                   always read correctly regardless.
    """
    level = paragraph.level or 0
    p_elm = paragraph._p  # read-only access to the underlying <a:p>; does
    # NOT use paragraph._pPr, which calls get_or_add_pPr() and would
    # mutate the presentation just to inspect it.
    pPr = p_elm.find(qn("a:pPr"))
    explicit = _bullet_from_properties(pPr, level)
    if explicit is not None:
        return explicit
    if _is_body_like_placeholder(shape):
        inherited = _resolve_inherited_bullet(shape, level)
        if inherited is not None:
            return inherited
    # A plain textbox/freeform with no bullet markup is prose, not a list.
    return {"is_bullet": False, "ordered": False, "level": level, "char": None}


def _render_text_frame(shape, is_title: bool = False) -> str:
    lines = []
    for para in shape.text_frame.paragraphs:
        text = "".join(run.text for run in para.runs).strip()
        if not text:
            continue
        if is_title:
            lines.append(text)
            continue
        info = _get_bullet_info(para, shape)
        indent = "  " * info["level"]
        if info["is_bullet"]:
            marker = "1." if info["ordered"] else "-"
            lines.append(f"{indent}{marker} {text}")
        else:
            lines.append(f"{indent}{text}" if info["level"] else text)
    return "\n".join(lines)


def _render_chart(chart) -> str:
    lines = []
    if chart.has_title:
        try:
            title = chart.chart_title.text_frame.text.strip()
            if title:
                lines.append(f"**Chart: {title}**")
        except Exception:
            pass
    if not lines:
        lines.append("**Chart**")
    try:
        for plot in chart.plots:
            categories = list(plot.categories)
            for series in plot.series:
                values = list(series.values)
                pairs = ", ".join(f"{c}: {v}" for c, v in zip(categories, values))
                lines.append(f"- {series.name}: {pairs}")
    except Exception:
        lines.append("<!-- chart data could not be read -->")
    return "\n".join(lines)


def _render_table(table) -> tuple:
    """Returns (markdown, flattened_grid, cells_or_None). `cells` carries
    explicit row/col/rowspan/colspan for merge-anchor cells so
    table_export.py can render standalone tables/*.html without flattening
    the span geometry away; `flattened_grid` (each cell's text, with
    spanned cells repeating "") is what standalone tables/*.csv uses."""
    n_rows = len(table.rows)
    n_cols = len(table.columns)
    grid = [["" for _ in range(n_cols)] for _ in range(n_rows)]
    cells = []
    any_span = False

    html = ["<table>"]
    for r in range(n_rows):
        html.append("<tr>")
        for c in range(n_cols):
            cell = table.cell(r, c)
            grid[r][c] = cell.text
            if cell.is_spanned:
                continue
            attrs = ""
            rowspan = colspan = 1
            if cell.is_merge_origin:
                any_span = True
                if cell.span_height > 1:
                    rowspan = cell.span_height
                    attrs += f' rowspan="{rowspan}"'
                if cell.span_width > 1:
                    colspan = cell.span_width
                    attrs += f' colspan="{colspan}"'
            html.append(f"<td{attrs}>{_html_escape(cell.text)}</td>")
            cells.append({"row": r, "col": c, "value": cell.text,
                          "rowspan": rowspan, "colspan": colspan})
        html.append("</tr>")
    html.append("</table>")

    if not any_span:
        lines = ["| " + " | ".join(_esc(grid[0][c]) for c in range(n_cols)) + " |",
                 "| " + " | ".join(["---"] * n_cols) + " |"]
        for r in range(1, n_rows):
            lines.append("| " + " | ".join(_esc(grid[r][c]) for c in range(n_cols)) + " |")
        return "\n".join(lines), grid, None

    return "\n".join(html), grid, cells


def _html_escape(s: str) -> str:
    return s.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")


def _esc(v) -> str:
    if v is None:
        return ""
    return str(v).replace("|", "\\|").replace("\n", "<br>")
