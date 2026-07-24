#!/usr/bin/env python3
"""Run runtime smoke qualification for the Agent Skill package profile."""
from __future__ import annotations
import argparse, json, os, shutil, subprocess, sys, zipfile
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]

def run(command, cwd, results, name):
    completed = subprocess.run(command, cwd=cwd, text=True, capture_output=True)
    results[name] = {"status": "passed" if completed.returncode == 0 else "failed", "command": command,
                     "stdout": completed.stdout[-2000:], "stderr": completed.stderr[-2000:]}
    if completed.returncode:
        raise RuntimeError(f"{name} failed")
    return completed

def safe_extract(archive: Path, extract_root: Path) -> Path:
    root = extract_root.resolve()
    with zipfile.ZipFile(archive) as zipped:
        for info in zipped.infolist():
            if info.is_dir():
                continue
            destination = (root / info.filename).resolve()
            try:
                destination.relative_to(root)
            except ValueError as exc:
                raise ValueError(f"unsafe extraction path: {info.filename!r}") from exc
            destination.parent.mkdir(parents=True, exist_ok=True)
            with zipped.open(info) as source, destination.open("wb") as target:
                shutil.copyfileobj(source, target)
    skill_dir = root / "universal-file-to-markdown"
    if not skill_dir.is_dir():
        raise RuntimeError(f"Expected extracted skill directory {skill_dir} not found")
    return skill_dir

def make_venv(interpreter: str, env_dir: Path, results: dict) -> Path:
    run([interpreter, "-m", "venv", "--clear", str(env_dir)], ROOT, results, "create_venv")
    python = env_dir / ("Scripts/python.exe" if os.name == "nt" else "bin/python")
    if not python.is_file():
        raise RuntimeError("requested interpreter did not create a usable virtual environment")
    version = run([str(python), "--version"], ROOT, results, "effective_python")
    results["effective_python"]["version"] = version.stdout.strip() or version.stderr.strip()
    return python

def make_fixtures(directory: Path, python: Path) -> None:
    script = '''from pathlib import Path
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from reportlab.pdfgen.canvas import Canvas
from PIL import Image, ImageDraw
p=Path(".")
d=Document(); d.add_paragraph("Package DOCX"); d.save(p/"sample.docx")
w=Workbook(); w.active["A1"]="Package XLSX"; w.save(p/"sample.xlsx")
r=Presentation(); r.slides.add_slide(r.slide_layouts[1]).shapes.title.text="Package PPTX"; r.save(p/"sample.pptx")
c=Canvas(str(p/"sample.pdf")); c.drawString(72,720,"Package PDF"); c.save()
(p/"sample.csv").write_bytes("name,city\\n中文,台北\\n".encode("big5"))
(p/"sample.json").write_text('{"message":"你好，package"}', encoding="utf-8")
(p/"sample.html").write_text("<h1>Package HTML</h1><p>ready</p>", encoding="utf-8")
image=Image.new("RGB",(300,100),"white"); ImageDraw.Draw(image).text((20,30),"PACKAGE OCR",fill="black"); image.save(p/"sample.png")
'''
    subprocess.run([str(python), "-c", script], cwd=directory, check=True)

def qualify_agent_skill(archive: Path, output: Path, requested_python: str) -> dict:
    archive = archive.resolve()
    output = output.resolve()
    output.mkdir(parents=True, exist_ok=True)
    results = {"archive": archive.name, "package_profile": "agent-skill", "status": "failed", "steps": {}}
    try:
        extract_root = output / "extract"
        shutil.rmtree(extract_root, ignore_errors=True)
        package = safe_extract(archive, extract_root)

        # Confirm tests and repo-only docs and manifests are excluded
        if (package / "tests").exists():
            raise RuntimeError("Agent Skill package must not contain tests/")
        if (package / "docs").exists():
            raise RuntimeError("Agent Skill package must not contain docs/")
        if (package / "package-manifest.json").exists():
            raise RuntimeError("Agent Skill package must not contain package-manifest.json")
        if (package / "package-manifests").exists():
            raise RuntimeError("Agent Skill package must not contain package-manifests/")

        # Validate package structurally inside isolated extracted folder without repository manifests
        run([sys.executable, str(package / "scripts/validate_skill_package.py"), str(archive),
             "--profile", "agent-skill"], package, results["steps"], "isolated_package_validation")

        env_dir = output / "venv"
        python = make_venv(requested_python, env_dir, results["steps"])
        run([str(python), "-m", "pip", "install", "--upgrade", "pip", "setuptools", "wheel"], package, results["steps"], "pip_bootstrap")
        run([str(python), "-m", "pip", "install", "-r", "requirements.txt"], package, results["steps"], "requirements_install")
        run([str(python), "scripts/capability_probe.py", "--json"], package, results["steps"], "capability_probe")

        # Run 8 representative conversions and bundle validations
        fixtures = output / "fixtures"
        fixtures.mkdir(parents=True, exist_ok=True)
        make_fixtures(fixtures, python)
        for extension in ("docx", "xlsx", "pptx", "pdf", "png", "csv", "json", "html"):
            bundle = output / "bundles" / extension
            run([str(python), str(package / "scripts/router.py"), str(fixtures / f"sample.{extension}"), "--output", str(bundle)], package, results["steps"], f"convert_{extension}")
            run([str(python), str(package / "scripts/validate_bundle.py"), str(bundle)], package, results["steps"], f"bundle_{extension}")
        results["status"] = "passed"
    except Exception as exc:
        results["error"] = str(exc)
    return results

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--archive", type=Path, required=True)
    parser.add_argument("--output", type=Path, default=Path(".qualification/agent-skill-smoke"))
    parser.add_argument("--python", default=sys.executable, help="interpreter used to create the clean venv")
    args = parser.parse_args()
    results = qualify_agent_skill(args.archive, args.output, args.python)
    print(json.dumps(results, ensure_ascii=False, indent=2))
    return 0 if results["status"] == "passed" else 1

if __name__ == "__main__":
    raise SystemExit(main())
