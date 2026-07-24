#!/usr/bin/env python3
"""Run clean-package qualification without trusting unvalidated archive content."""
from __future__ import annotations
import argparse
import json
import os
import shutil
import subprocess
import sys
import zipfile
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]


def run(command, cwd, results, name):
    completed = subprocess.run(command, cwd=cwd, text=True, capture_output=True)
    results[name] = {"status": "passed" if completed.returncode == 0 else "failed", "command": command,
                     "stdout": completed.stdout[-2000:], "stderr": completed.stderr[-2000:]}
    if completed.returncode:
        raise RuntimeError(f"{name} failed")
    return completed


def safe_extract(archive: Path, extract_root: Path) -> None:
    """Extract only regular, containment-checked members after validation."""
    root = extract_root.resolve()
    with zipfile.ZipFile(archive) as zipped:
        for info in zipped.infolist():
            if info.is_dir():
                continue
            destination = (root / info.filename).resolve()
            try:
                destination.relative_to(root)
            except ValueError as exc:
                raise ValueError(f"unsafe extraction path: {info.filename!r}") from exc
            destination.parent.mkdir(parents=True, exist_ok=True)
            with zipped.open(info) as source, destination.open("wb") as target:
                shutil.copyfileobj(source, target)


def make_venv(interpreter: str, env_dir: Path, results: dict) -> Path:
    run([interpreter, "-m", "venv", "--clear", str(env_dir)], ROOT, results, "create_venv")
    python = env_dir / ("Scripts/python.exe" if os.name == "nt" else "bin/python")
    if not python.is_file():
        raise RuntimeError("requested interpreter did not create a usable virtual environment")
    version = run([str(python), "--version"], ROOT, results, "effective_python")
    results["effective_python"]["version"] = version.stdout.strip() or version.stderr.strip()
    return python


def make_fixtures(directory: Path, python: Path) -> None:
    script = '''from pathlib import Path
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from reportlab.pdfgen.canvas import Canvas
from PIL import Image, ImageDraw
p=Path(".")
d=Document(); d.add_paragraph("Package DOCX"); d.save(p/"sample.docx")
w=Workbook(); w.active["A1"]="Package XLSX"; w.save(p/"sample.xlsx")
r=Presentation(); r.slides.add_slide(r.slide_layouts[1]).shapes.title.text="Package PPTX"; r.save(p/"sample.pptx")
c=Canvas(str(p/"sample.pdf")); c.drawString(72,720,"Package PDF"); c.save()
(p/"sample.csv").write_bytes("name,city\\n中文,台北\\n".encode("big5"))
(p/"sample.json").write_text('{"message":"你好，package"}', encoding="utf-8")
(p/"sample.html").write_text("<h1>Package HTML</h1><p>ready</p>", encoding="utf-8")
image=Image.new("RGB",(300,100),"white"); ImageDraw.Draw(image).text((20,30),"PACKAGE OCR",fill="black"); image.save(p/"sample.png")
'''
    subprocess.run([str(python), "-c", script], cwd=directory, check=True)


def qualify_conversions(package: Path, python: Path, output: Path, results: dict) -> None:
    fixtures = output / "fixtures"
    fixtures.mkdir(parents=True, exist_ok=True)
    make_fixtures(fixtures, python)
    for extension in ("docx", "xlsx", "pptx", "pdf", "png", "csv", "json", "html"):
        bundle = output / "bundles" / extension
        run([str(python), str(package / "scripts/router.py"), str(fixtures / f"sample.{extension}"), "--output", str(bundle)], package, results, f"convert_{extension}")
        run([str(python), str(package / "scripts/validate_bundle.py"), str(bundle)], package, results, f"bundle_{extension}")

def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--archive", type=Path, required=True)
    parser.add_argument("--output", type=Path, default=Path(".qualification/package-smoke"))
    parser.add_argument("--python", default=sys.executable, help="interpreter used to create the clean venv")
    parser.add_argument("--source-manifest", type=Path, default=ROOT / "package-manifest.json",
                        help="trusted local allowlist; never read from the archive before validation")
    args = parser.parse_args()
    archive, output = args.archive.resolve(), args.output.resolve()
    output.mkdir(parents=True, exist_ok=True)
    results = {"archive": archive.name, "requested_python": args.python, "status": "failed", "steps": {}}
    try:
        # The trusted manifest and local validator establish trust before any write from ZIP content.
        run([sys.executable, str(ROOT / "scripts/validate_skill_package.py"), str(archive),
             "--source-manifest", str(args.source_manifest.resolve())], ROOT, results["steps"], "archive_validation")
        extract_root = output / "extract"
        shutil.rmtree(extract_root, ignore_errors=True)
        safe_extract(archive, extract_root)
        package = next(extract_root.iterdir())
        env_dir = output / "venv"
        python = make_venv(args.python, env_dir, results["steps"])
        run([str(python), "-m", "pip", "install", "--upgrade", "pip", "setuptools", "wheel"], package, results["steps"], "pip_bootstrap")
        run([str(python), "-m", "pip", "install", "-r", "requirements.txt"], package, results["steps"], "requirements_install")
        run([str(python), "scripts/capability_probe.py", "--json"], package, results["steps"], "capability_probe")
        qualify_conversions(package, python, output, results["steps"])
        results["status"] = "passed"
    except Exception as exc:
        results["error"] = str(exc)
    print(json.dumps(results, ensure_ascii=False, indent=2))
    return 0 if results["status"] == "passed" else 1

if __name__ == "__main__":
    raise SystemExit(main())
