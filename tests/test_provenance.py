"""Phase 2 source-locator and chunk-provenance regression coverage."""
import os
import sys

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "scripts"))

from chunker import build_chunk_provenance
from validate_bundle import _valid_a1_range, _validate_source_locator


def _errors(locator):
    errors = []
    _validate_source_locator(locator, errors)
    return errors


def test_xlsx_chunk_inherits_sheet_cell_range_and_table_reference():
    element = {"table_id": "table-s001-b001", "locator_precision": "exact",
               "source_locator": {"format": "xlsx", "sheet_name": "Sheet1", "cell_range": "A1:G9"}}
    provenance = build_chunk_provenance([element], {})
    assert provenance["table_ids"] == ["table-s001-b001"]
    assert provenance["source_locator"]["cell_range"] == "A1:G9"
    assert provenance["locator_precision"] == "exact"


def test_pptx_multi_shape_and_multi_element_provenance():
    elements = [
        {"locator_precision": "exact", "source_locator": {"format": "pptx", "slide_number": 1, "shape_id": 5}},
        {"locator_precision": "exact", "source_locator": {"format": "pptx", "slide_number": 1, "shape_id": 7}},
    ]
    provenance = build_chunk_provenance(elements, {})
    assert provenance["locator_precision"] == "range"
    assert provenance["source_locator"]["shape_ids"] == [5, 7]


def test_pdf_ocr_docx_csv_json_and_eml_precision_contracts():
    assert build_chunk_provenance([{"locator_precision": "derived", "source_locator": {"format": "pdf", "page_start": 1, "page_end": 1}}], {})["locator_precision"] == "derived"
    assert build_chunk_provenance([{"locator_precision": "range", "source_locator": {"format": "docx", "section_index": 0, "element_start": 1, "element_end": 2}}], {})["locator_precision"] == "range"
    assert not _errors({"format": "csv", "row_start": 1, "row_end": 2})
    assert not _errors({"format": "json", "json_path": "$.items[0]"})
    assert not _errors({"format": "eml", "mime_part": "1.2", "section": "attachment"})


def test_invalid_locator_fixtures_are_rejected_with_stable_codes():
    assert _errors({"format": "xlsx", "sheet_name": "Sheet1", "cell_range": "B2:A1"}) == ["INVALID_XLSX_CELL_RANGE"]
    assert "INVALID_PDF_PAGE_RANGE" in _errors({"format": "pdf", "page_start": 2, "page_end": 1})
    assert "INVALID_PDF_BBOX" in _errors({"format": "pdf", "page_start": 1, "page_end": 1, "bboxes": [[2, 1, 0, 4]]})
    assert "INVALID_PPTX_SLIDE_NUMBER" in _errors({"format": "pptx", "slide_number": 0})
    assert "INVALID_CSV_ROW_RANGE" in _errors({"format": "csv", "row_start": 2, "row_end": 1})
    assert "INVALID_JSON_PATH" in _errors({"format": "json", "json_path": "items[0]"})
    assert _valid_a1_range("$A$1:$G$9")


def _csv_bundle(tmp_path):
    import router
    source = tmp_path / "source.csv"
    source.write_text("name,value\nA,1\n", encoding="utf-8")
    output = tmp_path / "bundle"
    assert router.convert(str(source), str(output))["status"] == "passed"
    return output


def test_current_router_output_has_precision_bounded_chunks_and_resolvable_references(tmp_path):
    import json
    from run_provenance_regression import metrics
    result = metrics(_csv_bundle(tmp_path))
    assert result["chunks_without_precision"] == result["chunks_over_2000"] == 0
    assert result["unresolved_element_refs"] == result["unresolved_table_refs"] == 0
    assert result["formats"]["xlsx"]["status"] == "not_applicable"
    assert result["formats"]["pptx"]["coverage"] is None
    assert result["formats"]["pdf"]["eligible_chunks"] == 0


def test_bundle_rejects_missing_references_conflicts_and_invalid_element_locator(tmp_path):
    import json
    from validate_bundle import validate_bundle
    bundle = _csv_bundle(tmp_path)
    chunks_path = bundle / "chunks.jsonl"
    chunk = json.loads(chunks_path.read_text(encoding="utf-8"))
    chunk["element_ids"] = ["missing"]
    chunk["table_ids"] = ["missing-table"]
    chunk["source_locators"] = [chunk["source_locator"]]
    chunks_path.write_text(json.dumps(chunk) + "\n", encoding="utf-8")
    errors = validate_bundle(str(bundle))["errors"]
    assert any(x.startswith("CHUNK_ELEMENT_REFERENCE_MISSING") for x in errors)
    assert any(x.startswith("CHUNK_TABLE_REFERENCE_MISSING") for x in errors)
    assert "CHUNK_LOCATOR_CONFLICT" in errors
    document_path = bundle / "document.json"
    document = json.loads(document_path.read_text(encoding="utf-8"))
    document["elements"][1]["locator_precision"] = "exact"
    document["elements"][1]["source_locator"] = {"format": "xlsx", "sheet_name": "Sheet", "cell_range": "B2:A1"}
    document_path.write_text(json.dumps(document), encoding="utf-8")
    assert "INVALID_XLSX_CELL_RANGE" in validate_bundle(str(bundle))["errors"]


def test_legacy_v160_chunk_without_provenance_remains_accepted(tmp_path):
    import json
    from validate_bundle import validate_bundle
    bundle = _csv_bundle(tmp_path)
    chunks_path = bundle / "chunks.jsonl"
    chunk = json.loads(chunks_path.read_text(encoding="utf-8"))
    for field in ("locator_precision", "source_locator", "source_locators", "table_ids"):
        chunk.pop(field, None)
    chunks_path.write_text(json.dumps(chunk) + "\n", encoding="utf-8")
    assert validate_bundle(str(bundle))["status"] == "passed"


def test_xlsx_overlapping_and_adjacent_ranges_merge():
    from chunker import merge_source_locators
    merged, precision = merge_source_locators([
        {"format": "xlsx", "sheet_name": "S", "cell_range": "A1:A2"},
        {"format": "xlsx", "sheet_name": "S", "cell_range": "A2:A3"},
    ])
    assert precision == "range" and merged["cell_range"] == "A1:A3"


def test_xlsx_disjoint_ranges_remain_multiple_locators():
    from chunker import merge_source_locators
    merged, precision = merge_source_locators([
        {"format": "xlsx", "sheet_name": "S", "cell_range": "A1:A1"},
        {"format": "xlsx", "sheet_name": "S", "cell_range": "Z100:Z100"},
    ])
    assert precision == "range" and isinstance(merged, list) and len(merged) == 2


def test_heading_hierarchy_is_preserved_and_heading_is_not_coalesced():
    from chunker import build_chunks
    elements = [
        {"id": "h1", "type": "heading", "parent_id": "root", "child_ids": [], "content": "# Title", "heading_path": [], "source_locator": {"format": "json", "json_path": "$.title"}, "locator_precision": "exact"},
        {"id": "p1", "type": "paragraph", "parent_id": "h1", "child_ids": [], "content": "Body", "heading_path": [], "source_locator": {"format": "json", "json_path": "$.body"}, "locator_precision": "exact"},
        {"id": "h2", "type": "heading", "parent_id": "h1", "child_ids": [], "content": "## Nested", "heading_path": ["Title"], "source_locator": {"format": "json", "json_path": "$.nested"}, "locator_precision": "exact"},
    ]
    chunks = build_chunks("", elements, "a" * 64)
    assert chunks[0]["element_ids"] == ["h1"] and chunks[0]["heading_path"] == ["Title"]
    assert chunks[1]["element_ids"] == ["p1"] and chunks[1]["heading_path"] == ["Title"]
    assert chunks[2]["element_ids"] == ["h2"] and chunks[2]["heading_path"] == ["Title", "Nested"]


def _generated_bundle(tmp_path, kind):
    import router
    output = tmp_path / f"{kind}-bundle"
    source = tmp_path / f"sample.{kind}"
    if kind == "xlsx":
        import openpyxl
        wb = openpyxl.Workbook(); ws = wb.active; ws.title = "Data"; ws.append(["A", "B"]); ws.append([1, 2]); wb.save(source)
    elif kind == "pptx":
        from pptx import Presentation
        from pptx.util import Inches
        presentation = Presentation(); slide = presentation.slides.add_slide(presentation.slide_layouts[6]); slide.shapes.add_textbox(Inches(1), Inches(1), Inches(2), Inches(1)).text_frame.text = "Slide text"; presentation.save(source)
    else:
        from reportlab.pdfgen.canvas import Canvas
        canvas = Canvas(str(source)); canvas.drawString(72, 720, "Digital PDF text"); canvas.save()
    assert router.convert(str(source), str(output))["status"] == "passed"
    return output


def test_generated_xlsx_pptx_pdf_metrics_are_non_vacuous(tmp_path):
    from run_provenance_regression import metrics
    for kind in ("xlsx", "pptx", "pdf"):
        result = metrics(_generated_bundle(tmp_path, kind))["formats"][kind]
        assert result["eligible_chunks"] > 0
        assert result["located_chunks"] == result["eligible_chunks"]
        assert result["coverage"] == 1.0 and result["status"] == "passed"
