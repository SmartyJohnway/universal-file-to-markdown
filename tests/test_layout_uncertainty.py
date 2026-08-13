from pathlib import Path
import sys

ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "scripts"))

from pdf_converter import _analyze_digital_layout
from pptx_converter import _analyze_slide_layout, convert_pptx


def test_pdf_two_column_geometry_is_flagged_for_deterministic_ordering():
    blocks = [
        ((40, 100, 220, 140), "A1"),
        ((40, 180, 220, 220), "A2"),
        ((360, 105, 540, 145), "B1"),
        ((360, 185, 540, 225), "B2"),
    ]
    result = _analyze_digital_layout(blocks, 600)
    assert result["multi_column_detected"] is True
    assert result["column_count"] == 2
    assert result["left_block_count"] == result["right_block_count"] == 2


def test_pdf_single_column_and_full_width_title_are_not_flagged():
    blocks = [
        ((30, 20, 570, 70), "Full width title"),
        ((40, 100, 560, 145), "Paragraph 1"),
        ((40, 170, 560, 215), "Paragraph 2"),
        ((40, 240, 560, 285), "Paragraph 3"),
    ]
    assert _analyze_digital_layout(blocks, 600)["multi_column_detected"] is False


def test_pptx_side_by_side_flows_are_flagged():
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    left = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4), Inches(3))
    left.text = "Left flow"
    right = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(4), Inches(3))
    right.text = "Right flow"
    result = _analyze_slide_layout(slide.shapes, prs.slide_width, prs.slide_height)
    assert result["visual_flow_ambiguous"] is True
    assert "independent_side_by_side_flows" in result["signals"]


def test_pptx_typical_title_and_body_are_not_flagged():
    from pptx import Presentation

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Title"
    slide.placeholders[1].text = "Body"
    result = _analyze_slide_layout(slide.shapes, prs.slide_width, prs.slide_height)
    assert result["visual_flow_ambiguous"] is False


def test_pptx_converter_emits_truthful_visual_flow_warning(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    source = tmp_path / "multi-flow.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for left, text in ((0.5, "Left flow"), (5.0, "Right flow")):
        box = slide.shapes.add_textbox(Inches(left), Inches(1.5), Inches(4), Inches(3))
        box.text = text
    prs.save(source)

    result = convert_pptx(str(source))
    assert result["report"]["status"] == "passed_with_warnings"
    assert [warning["code"] for warning in result["report"]["warnings"]] == ["VISUAL_FLOW_AMBIGUOUS"]
