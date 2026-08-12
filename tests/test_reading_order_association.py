import json
from pathlib import Path
import sys


ROOT = Path(__file__).resolve().parents[1]
sys.path.insert(0, str(ROOT / "scripts"))

from pdf_converter import (
    _analyze_digital_layout,
    _associate_pdf_page_captions,
    _plan_pdf_page_items,
    convert_pdf,
)
from pptx_converter import _plan_slide_reading_order, convert_pptx
from validate_bundle import validate_bundle
import router


def _two_column_blocks():
    return [
        ((20, 20, 580, 60), "Report Title"),
        ((40, 100, 220, 140), "A1"),
        ((360, 105, 540, 145), "B1"),
        ((40, 180, 220, 220), "A2"),
        ((360, 185, 540, 225), "B2"),
        ((20, 700, 580, 740), "Page Footer"),
    ]


def test_pdf_xycut_plan_orders_columns_and_preserves_spanning_blocks():
    blocks = _two_column_blocks()
    layout = _analyze_digital_layout(blocks, 600)
    plan = _plan_pdf_page_items(blocks, [], layout, 600, 800, "page-0001")

    assert [item["content"] for item in plan] == [
        "Report Title", "A1", "A2", "B1", "B2", "Page Footer"
    ]
    assert [item["layout"]["reading_order"] for item in plan] == list(range(1, 7))
    assert [item["layout"]["column_index"] for item in plan[1:5]] == [1, 1, 2, 2]
    assert all(item["layout"]["order_method"] == "deterministic_xycut_v1" for item in plan)


def test_pdf_single_column_table_is_inserted_by_bbox_not_appended():
    blocks = [
        ((40, 80, 560, 120), "Before table"),
        ((40, 300, 560, 340), "After table"),
    ]
    layout = _analyze_digital_layout(blocks, 600)
    table = {
        "id": "table-p0001-0001",
        "element_id": "page-0001-table-001",
        "rows": [["A", "B"]],
        "markdown": "| A | B |",
        "bbox": (40, 180, 560, 240),
    }
    plan = _plan_pdf_page_items(blocks, [table], layout, 600, 800, "page-0001")
    assert [item["content"] for item in plan] == [
        "Before table", "| A | B |", "After table"
    ]


def test_pdf_side_note_is_not_promoted_to_second_column_without_support():
    blocks = [
        ((40, 100, 430, 150), "Main 1"),
        ((40, 180, 430, 230), "Main 2"),
        ((470, 140, 560, 200), "One side note"),
        ((40, 260, 430, 310), "Main 3"),
    ]
    assert _analyze_digital_layout(blocks, 600)["multi_column_detected"] is False


def test_pdf_overlap_reduces_order_confidence():
    blocks = [
        ((40, 100, 220, 160), "A1"),
        ((40, 130, 220, 190), "A2 overlapping"),
        ((360, 100, 540, 150), "B1"),
        ((360, 180, 540, 230), "B2"),
    ]
    layout = _analyze_digital_layout(blocks, 600)
    assert layout["multi_column_detected"] is True
    assert layout["overlapping_block_pair_count"] == 1
    plan = _plan_pdf_page_items(blocks, [], layout, 600, 800, "page-0001")
    assert {item["layout"]["order_confidence"] for item in plan} == {0.65}


def test_pdf_caption_association_requires_prefix_proximity_and_alignment():
    caption = {
        "id": "caption", "type": "paragraph", "content": "Table 1: Results",
        "source_locator": {"bbox": [40, 100, 300, 125]}, "properties": {},
    }
    table = {
        "id": "table", "type": "table", "content": "| A |",
        "source_locator": {"bbox": [40, 130, 560, 260]}, "properties": {},
    }
    unrelated = {
        "id": "prose", "type": "paragraph", "content": "Table manners matter",
        "source_locator": {"bbox": [40, 300, 300, 325]}, "properties": {},
    }
    _associate_pdf_page_captions([caption, table, unrelated], 800)
    assert caption["type"] == "caption"
    assert caption["properties"]["associations"][0]["relation"] == "caption_of"
    assert caption["properties"]["associations"][0]["target_id"] == "table"
    assert "associations" not in unrelated["properties"]


def _make_two_column_pdf(tmp_path):
    import fitz

    source = tmp_path / "two-columns.pdf"
    document = fitz.open()
    page = document.new_page(width=600, height=800)
    for bbox, text in _two_column_blocks():
        page.insert_textbox(fitz.Rect(*bbox), text, fontsize=11)
    document.save(source)
    document.close()
    return source


def test_real_pdf_projection_and_canonical_share_column_major_order(tmp_path):
    source = _make_two_column_pdf(tmp_path)

    result = convert_pdf(str(source))
    positions = [result["markdown"].index(text) for text in
                 ("Report Title", "A1", "A2", "B1", "B2", "Page Footer")]
    assert positions == sorted(positions)
    children = [element for element in result["elements"]
                if element.get("parent_id") == "page-0001"]
    contents = [element["content"] for element in children]
    assert contents == ["Report Title", "A1", "A2", "B1", "B2", "Page Footer"]


def test_pdf_bundle_chunks_follow_canonical_order_and_rerun_deterministically(tmp_path):
    source = _make_two_column_pdf(tmp_path)
    outputs = [tmp_path / "out-1", tmp_path / "out-2"]
    for output in outputs:
        assert router.convert(str(source), str(output))["bundle_validation_status"] == "passed"

    first_document = json.loads((outputs[0] / "document.json").read_text(encoding="utf-8"))
    child_ids = [element["id"] for element in first_document["elements"]
                 if element.get("parent_id") == "page-0001"]
    chunks = [json.loads(line) for line in
              (outputs[0] / "chunks.jsonl").read_text(encoding="utf-8").splitlines()]
    chunk_element_ids = [element_id for chunk in chunks for element_id in chunk["element_ids"]]
    assert chunk_element_ids == child_ids
    for name in ("document.md", "document.json", "chunks.jsonl"):
        assert (outputs[0] / name).read_bytes() == (outputs[1] / name).read_bytes()


def test_real_pdf_table_is_inserted_between_surrounding_text(tmp_path):
    from reportlab.pdfgen.canvas import Canvas

    source = tmp_path / "table-order.pdf"
    canvas = Canvas(str(source), pagesize=(600, 800))
    canvas.drawString(50, 750, "Before table")
    xs, ys = [50, 250, 450], [650, 620, 590]
    for x in xs:
        canvas.line(x, ys[-1], x, ys[0])
    for y in ys:
        canvas.line(xs[0], y, xs[-1], y)
    canvas.drawString(60, 630, "Header A")
    canvas.drawString(260, 630, "Header B")
    canvas.drawString(60, 600, "Value A")
    canvas.drawString(260, 600, "Value B")
    canvas.drawString(50, 500, "After table")
    canvas.save()

    result = convert_pdf(str(source))
    assert len(result["tables"]) == 1
    positions = [result["markdown"].index(value)
                 for value in ("Before table", "Header A", "After table")]
    assert positions == sorted(positions)
    assert [element["type"] for element in result["elements"]] == [
        "page", "heading", "table", "heading"
    ]


def _add_textbox(slide, inches, left, top, width, height, text):
    box = slide.shapes.add_textbox(
        inches(left), inches(top), inches(width), inches(height)
    )
    box.text = text
    return box


def test_pptx_plan_uses_column_major_order_with_geometry_tie_breakers():
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches, 0.3, 0.2, 9.3, 0.5, "Title")
    _add_textbox(slide, Inches, 0.5, 1.2, 4.0, 0.6, "A1")
    _add_textbox(slide, Inches, 5.0, 1.2, 4.0, 0.6, "B1")
    _add_textbox(slide, Inches, 0.5, 2.5, 4.0, 0.6, "A2")
    _add_textbox(slide, Inches, 5.0, 2.5, 4.0, 0.6, "B2")
    _add_textbox(slide, Inches, 0.3, 6.8, 9.3, 0.3, "Footer")

    ordered, plan, analysis = _plan_slide_reading_order(
        slide.shapes, prs.slide_width, prs.slide_height
    )
    assert [shape.text for shape in ordered if getattr(shape, "has_text_frame", False)] == [
        "Title", "A1", "A2", "B1", "B2", "Footer"
    ]
    assert analysis["column_count"] == 2
    assert sorted(value["reading_order"] for value in plan.values()) == list(range(1, 7))


def test_pptx_converter_projection_and_elements_share_plan(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    source = tmp_path / "multi-flow.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    for left, top, text in (
        (0.5, 1.2, "A1"), (5.0, 1.2, "B1"),
        (0.5, 2.5, "A2"), (5.0, 2.5, "B2"),
    ):
        _add_textbox(slide, Inches, left, top, 4.0, 0.6, text)
    prs.save(source)

    result = convert_pptx(str(source))
    positions = [result["markdown"].index(text) for text in ("A1", "A2", "B1", "B2")]
    assert positions == sorted(positions)
    children = [element for element in result["elements"]
                if element.get("parent_id") == "slide-0001"]
    assert [element["content"] for element in children] == ["A1", "A2", "B1", "B2"]
    assert all((element.get("properties") or {}).get("layout") for element in children)


def test_pptx_table_caption_and_speaker_note_edges_validate_in_bundle(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    source = tmp_path / "associations.pptx"
    output = tmp_path / "out"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    caption = _add_textbox(slide, Inches, 1.0, 0.7, 5.0, 0.4, "Table 1: Results")
    table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1.2), Inches(5), Inches(2))
    table_shape.table.cell(0, 0).text = "A"
    table_shape.table.cell(0, 1).text = "B"
    slide.notes_slide.notes_text_frame.text = "Explain the result"
    prs.save(source)

    report = router.convert(str(source), str(output))
    assert report["status"] in {"passed", "passed_with_warnings"}
    canonical = json.loads((output / "document.json").read_text(encoding="utf-8"))
    by_type = {}
    for element in canonical["elements"]:
        by_type.setdefault(element["type"], []).append(element)
    caption_element = by_type["caption"][0]
    table_element = by_type["table"][0]
    note_element = by_type["speaker_note"][0]
    assert caption_element["properties"]["associations"][0]["target_id"] == table_element["id"]
    assert note_element["properties"]["associations"][0] == {
        "relation": "note_for",
        "target_id": "slide-0001",
        "confidence": 1.0,
        "evidence": ["ooxml_notes_relationship"],
        "method": "deterministic_rule_v1",
    }
    assert report["bundle_validation_status"] == "passed"


def test_bundle_validator_rejects_duplicate_sibling_reading_order(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    source = tmp_path / "duplicate-order.pptx"
    output = tmp_path / "out"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches, 1.0, 1.0, 3.0, 0.5, "First")
    _add_textbox(slide, Inches, 1.0, 2.0, 3.0, 0.5, "Second")
    prs.save(source)
    assert router.convert(str(source), str(output))["bundle_validation_status"] == "passed"

    path = output / "document.json"
    canonical = json.loads(path.read_text(encoding="utf-8"))
    children = [element for element in canonical["elements"]
                if element.get("parent_id") == "slide-0001"]
    children[1]["properties"]["layout"]["reading_order"] = (
        children[0]["properties"]["layout"]["reading_order"]
    )
    path.write_text(json.dumps(canonical), encoding="utf-8")
    result = validate_bundle(str(output))
    assert result["status"] == "failed"
    assert any("LAYOUT_READING_ORDER_DUPLICATE" in error for error in result["errors"])


def test_bundle_validator_rejects_missing_association_target(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    source = tmp_path / "bad-edge.pptx"
    output = tmp_path / "out"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_textbox(slide, Inches, 1.0, 0.7, 5.0, 0.4, "Table 1: Results")
    slide.shapes.add_table(2, 2, Inches(1), Inches(1.2), Inches(5), Inches(2))
    prs.save(source)
    assert router.convert(str(source), str(output))["bundle_validation_status"] == "passed"

    path = output / "document.json"
    canonical = json.loads(path.read_text(encoding="utf-8"))
    caption = next(element for element in canonical["elements"]
                   if element["type"] == "caption")
    caption["properties"]["associations"][0]["target_id"] = "missing-element"
    path.write_text(json.dumps(canonical), encoding="utf-8")
    result = validate_bundle(str(output))
    assert result["status"] == "failed"
    assert any("ASSOCIATION_TARGET_MISSING" in error for error in result["errors"])
