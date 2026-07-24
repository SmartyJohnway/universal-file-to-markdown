"""Focused production-path regression for explicit target-table reason codes."""

import json
import subprocess
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
SCRIPTS = ROOT / "scripts"


def _run(command: list[str]) -> subprocess.CompletedProcess[str]:
    return subprocess.run(command, cwd=ROOT, text=True, capture_output=True)


def _load_json(path: Path) -> dict:
    return json.loads(path.read_text(encoding="utf-8"))


def _table_id(bundle: Path) -> str:
    tables_dir = bundle / "tables"
    files = sorted(p for p in tables_dir.glob("*.json") if p.name != "index.json")
    if not files:
        raise AssertionError(f"No table files found in {tables_dir}")
    return _load_json(files[0])["id"]


def test_plain_docx_explicit_target_table_reasons(tmp_path):
    from docx import Document

    source_dir = tmp_path / "sources" / "docx_plain"
    bundle_dir = tmp_path / "bundles" / "docx_plain"
    source_dir.mkdir(parents=True, exist_ok=True)
    bundle_dir.mkdir(parents=True, exist_ok=True)
    source = source_dir / "sample.docx"

    doc = Document()
    doc.add_paragraph("Plain table test")
    t = doc.add_table(rows=2, cols=2)
    t.cell(0, 0).text = "A1"
    t.cell(0, 1).text = "B1"
    t.cell(1, 0).text = "A2"
    t.cell(1, 1).text = "B2"
    doc.save(source)

    res = _run([sys.executable, str(SCRIPTS / "router.py"), str(source), "--output", str(bundle_dir)])
    assert res.returncode == 0, res.stderr

    tid = _table_id(bundle_dir)

    prep = _run(
        [
            sys.executable,
            str(SCRIPTS / "prepare_ai_review.py"),
            str(bundle_dir),
            "--force-user-request",
            "--target-table",
            tid,
        ]
    )
    assert prep.returncode == 0, prep.stderr

    req = _load_json(bundle_dir / "ai-review-request.json")
    assert "EXPLICIT_USER_REQUEST" in req["reason_codes"]
    assert "MERGED_TABLE_GEOMETRY_PRESENT" not in req["reason_codes"]
    assert "HTML_MERGED_TABLE_COMPLEX" not in req["reason_codes"]

    assert len(req["targets"]) == 1
    target = req["targets"][0]
    assert target["target_type"] == "table"
    assert target["target_id"] == tid
    assert "EXPLICIT_USER_REQUEST" in target["reason_codes"]
    assert "MERGED_TABLE_GEOMETRY_PRESENT" not in target["reason_codes"]
    assert "HTML_MERGED_TABLE_COMPLEX" not in target["reason_codes"]


def test_plain_xlsx_explicit_target_table_reasons(tmp_path):
    import openpyxl

    source_dir = tmp_path / "sources" / "xlsx_plain"
    bundle_dir = tmp_path / "bundles" / "xlsx_plain"
    source_dir.mkdir(parents=True, exist_ok=True)
    bundle_dir.mkdir(parents=True, exist_ok=True)
    source = source_dir / "sample.xlsx"

    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "A1"
    ws["B1"] = "B1"
    ws["A2"] = "A2"
    ws["B2"] = "B2"
    wb.save(source)

    res = _run([sys.executable, str(SCRIPTS / "router.py"), str(source), "--output", str(bundle_dir)])
    assert res.returncode == 0, res.stderr

    tid = _table_id(bundle_dir)

    prep = _run(
        [
            sys.executable,
            str(SCRIPTS / "prepare_ai_review.py"),
            str(bundle_dir),
            "--force-user-request",
            "--target-table",
            tid,
        ]
    )
    assert prep.returncode == 0, prep.stderr

    req = _load_json(bundle_dir / "ai-review-request.json")
    assert "EXPLICIT_USER_REQUEST" in req["reason_codes"]
    assert "MERGED_TABLE_GEOMETRY_PRESENT" not in req["reason_codes"]
    assert "HTML_MERGED_TABLE_COMPLEX" not in req["reason_codes"]

    assert len(req["targets"]) == 1
    target = req["targets"][0]
    assert target["target_type"] == "table"
    assert target["target_id"] == tid
    assert "EXPLICIT_USER_REQUEST" in target["reason_codes"]
    assert "MERGED_TABLE_GEOMETRY_PRESENT" not in target["reason_codes"]
    assert "HTML_MERGED_TABLE_COMPLEX" not in target["reason_codes"]


def test_plain_pptx_explicit_target_table_reasons(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    source_dir = tmp_path / "sources" / "pptx_plain"
    bundle_dir = tmp_path / "bundles" / "pptx_plain"
    source_dir.mkdir(parents=True, exist_ok=True)
    bundle_dir.mkdir(parents=True, exist_ok=True)
    source = source_dir / "sample.pptx"

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(5), Inches(2))
    table = shape.table
    table.cell(0, 0).text = "A1"
    table.cell(0, 1).text = "B1"
    table.cell(1, 0).text = "A2"
    table.cell(1, 1).text = "B2"
    prs.save(source)

    res = _run([sys.executable, str(SCRIPTS / "router.py"), str(source), "--output", str(bundle_dir)])
    assert res.returncode == 0, res.stderr

    tid = _table_id(bundle_dir)

    prep = _run(
        [
            sys.executable,
            str(SCRIPTS / "prepare_ai_review.py"),
            str(bundle_dir),
            "--force-user-request",
            "--target-table",
            tid,
        ]
    )
    assert prep.returncode == 0, prep.stderr

    req = _load_json(bundle_dir / "ai-review-request.json")
    assert "EXPLICIT_USER_REQUEST" in req["reason_codes"]
    assert "MERGED_TABLE_GEOMETRY_PRESENT" not in req["reason_codes"]
    assert "HTML_MERGED_TABLE_COMPLEX" not in req["reason_codes"]

    assert len(req["targets"]) == 1
    target = req["targets"][0]
    assert target["target_type"] == "table"
    assert target["target_id"] == tid
    assert "EXPLICIT_USER_REQUEST" in target["reason_codes"]
    assert "MERGED_TABLE_GEOMETRY_PRESENT" not in target["reason_codes"]
    assert "HTML_MERGED_TABLE_COMPLEX" not in target["reason_codes"]


def test_plain_html_explicit_target_table_reasons(tmp_path):
    source_dir = tmp_path / "sources" / "html_plain"
    bundle_dir = tmp_path / "bundles" / "html_plain"
    source_dir.mkdir(parents=True, exist_ok=True)
    bundle_dir.mkdir(parents=True, exist_ok=True)
    source = source_dir / "sample.html"

    source.write_text(
        "<!doctype html><html><head><title>Plain HTML</title></head><body>"
        "<table><tr><th>Header A</th><th>Header B</th></tr>"
        "<tr><td>Value A</td><td>Value B</td></tr></table>"
        "</body></html>",
        encoding="utf-8",
    )

    res = _run(
        [
            sys.executable,
            str(SCRIPTS / "router.py"),
            str(source),
            "--output",
            str(bundle_dir),
            "--source-url",
            "https://example.test/plain.html",
        ]
    )
    assert res.returncode == 0, res.stderr

    tid = _table_id(bundle_dir)

    prep = _run(
        [
            sys.executable,
            str(SCRIPTS / "prepare_ai_review.py"),
            str(bundle_dir),
            "--force-user-request",
            "--target-table",
            tid,
        ]
    )
    assert prep.returncode == 0, prep.stderr

    req = _load_json(bundle_dir / "ai-review-request.json")
    assert "EXPLICIT_USER_REQUEST" in req["reason_codes"]
    assert "MERGED_TABLE_GEOMETRY_PRESENT" not in req["reason_codes"]
    assert "HTML_MERGED_TABLE_COMPLEX" not in req["reason_codes"]

    assert len(req["targets"]) == 1
    target = req["targets"][0]
    assert target["target_type"] == "table"
    assert target["target_id"] == tid
    assert "EXPLICIT_USER_REQUEST" in target["reason_codes"]
    assert "MERGED_TABLE_GEOMETRY_PRESENT" not in target["reason_codes"]
    assert "HTML_MERGED_TABLE_COMPLEX" not in target["reason_codes"]
