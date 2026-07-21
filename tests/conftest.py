"""
conftest.py
Generates every test fixture programmatically at test time rather than
committing binary .xlsx/.docx/.pptx/.pdf files to the repo - keeps the
test suite self-contained, diff-friendly, and reproducible in any
environment without relying on checked-in binaries.
"""

import os
import sys
import zipfile
from xml.etree import ElementTree as ET

import pytest

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "scripts"))


@pytest.fixture
def tmp_out(tmp_path):
    return str(tmp_path / "out")


@pytest.fixture
def xlsx_merged(tmp_path):
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"] = "Header"
    ws.merge_cells("A1:B1")
    ws["A2"] = "X"
    ws["B2"] = "Y"
    path = str(tmp_path / "merged.xlsx")
    wb.save(path)
    return path


@pytest.fixture
def xlsx_formula_no_cache(tmp_path):
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A3"] = 5
    ws["B3"] = 7
    ws["C3"] = "=A3+B3"
    path = str(tmp_path / "formula.xlsx")
    wb.save(path)
    return path


@pytest.fixture
def docx_horizontal_merge(tmp_path):
    from docx import Document
    doc = Document()
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).merge(table.cell(0, 1))
    table.cell(0, 0).text = "Merged Header"
    table.cell(1, 0).text = "A"
    table.cell(1, 1).text = "B"
    path = str(tmp_path / "hmerge.docx")
    doc.save(path)
    return path


@pytest.fixture
def docx_vertical_merge(tmp_path):
    from docx import Document
    doc = Document()
    table = doc.add_table(rows=2, cols=2)
    a = table.cell(0, 0)
    a.merge(table.cell(1, 0))
    a.text = "VMerge"
    table.cell(0, 1).text = "X1"
    table.cell(1, 1).text = "X2"
    path = str(tmp_path / "vmerge.docx")
    doc.save(path)
    return path


@pytest.fixture
def docx_bold_italic(tmp_path):
    from docx import Document
    doc = Document()
    p = doc.add_paragraph()
    r1 = p.add_run("bold text")
    r1.bold = True
    p.add_run(" normal ")
    r3 = p.add_run("italic")
    r3.italic = True
    path = str(tmp_path / "runs.docx")
    doc.save(path)
    return path


@pytest.fixture
def pdf_short_digital(tmp_path):
    from reportlab.pdfgen import canvas
    path = str(tmp_path / "short.pdf")
    c = canvas.Canvas(path)
    c.drawString(100, 700, "Digital PDF Test 123")
    c.save()
    return path


@pytest.fixture
def csv_big5(tmp_path):
    text = "名稱,數量\n鋼管,10\n"
    path = str(tmp_path / "big5.csv")
    with open(path, "wb") as f:
        f.write(text.encode("big5"))
    return path


@pytest.fixture
def image_with_mock_ocr(tmp_path, monkeypatch):
    from PIL import Image
    import rapidocr_onnxruntime
    path = str(tmp_path / "scan.png")
    Image.new("RGB", (300, 100), "white").save(path)

    class FakeRapidOCR:
        def __call__(self, _image_bytes):
            return [
                ([[10, 10], [100, 10], [100, 30], [10, 30]], "Hello", 0.98),
                ([[110, 10], [200, 10], [200, 30], [110, 30]], "world", 0.97),
            ], None

    monkeypatch.setattr(rapidocr_onnxruntime, "RapidOCR", FakeRapidOCR)
    return path


@pytest.fixture
def xlsx_mislabeled_as_pdf(tmp_path, xlsx_merged):
    path = str(tmp_path / "mislabeled.pdf")
    with open(xlsx_merged, "rb") as src, open(path, "wb") as dst:
        dst.write(src.read())
    return path


@pytest.fixture
def xlsx_two_blocks(tmp_path):
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"], ws["B1"] = "Left", "Value"
    ws["A2"], ws["B2"] = "A", 1
    ws["A5"], ws["B5"] = "Second", "Value"
    ws["A6"], ws["B6"] = "B", 2
    path = str(tmp_path / "two_blocks.xlsx")
    wb.save(path)
    return path


@pytest.fixture
def pptx_merge_table(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(2))
    table = shape.table
    table.cell(0, 0).text = "Merged"
    table.cell(0, 0).merge(table.cell(0, 1))
    table.cell(1, 0).text = "A"
    table.cell(1, 1).text = "B"
    path = str(tmp_path / "merge.pptx")
    prs.save(path)
    return path


@pytest.fixture
def pptx_picture(tmp_path):
    from PIL import Image
    from pptx import Presentation
    from pptx.util import Inches
    image_path = str(tmp_path / "pixel.png")
    Image.new("RGB", (4, 4), (255, 0, 0)).save(image_path)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(image_path, Inches(1), Inches(1), Inches(1), Inches(1))
    path = str(tmp_path / "picture.pptx")
    prs.save(path)
    return path


@pytest.fixture
def pptx_level0_bullets(tmp_path):
    """A text box with two top-level (level 0) bulleted paragraphs and
    one explicit non-bullet paragraph (buNone), to regression-test that
    level 0 is no longer treated as "not a bullet"."""
    from pptx import Presentation
    from pptx.oxml.ns import qn
    from lxml import etree
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    box = slide.placeholders[1]
    tf = box.text_frame
    p0 = tf.paragraphs[0]
    p0.text = "First bullet point"
    p1 = tf.add_paragraph()
    p1.text = "Second bullet point"
    p2 = tf.add_paragraph()
    p2.text = "Not a bullet line"
    # p0/p1 inherit bullets from the body placeholder's layout/master.
    # p2 gets an explicit <a:buNone/> so it must NOT render as a bullet.
    pPr = p2._p.get_or_add_pPr()
    etree.SubElement(pPr, qn("a:buNone"))
    path = str(tmp_path / "bullets.pptx")
    prs.save(path)
    return path


@pytest.fixture
def pptx_plain_textbox(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    box.text_frame.text = "Plain paragraph one"
    box.text_frame.add_paragraph().text = "Plain paragraph two"
    path = str(tmp_path / "plain_textbox.pptx")
    prs.save(path)
    return path


@pytest.fixture
def pptx_explicit_bullets(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches
    from pptx.oxml.ns import qn
    from lxml import etree
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(3))
    tf = box.text_frame
    p0 = tf.paragraphs[0]
    p0.text = "Explicit char"
    etree.SubElement(p0._p.get_or_add_pPr(), qn("a:buChar"), {"char": "•"})
    p1 = tf.add_paragraph()
    p1.text = "Explicit number"
    etree.SubElement(p1._p.get_or_add_pPr(), qn("a:buAutoNum"), {"type": "arabicPeriod"})
    path = str(tmp_path / "explicit_bullets.pptx")
    prs.save(path)
    return path


@pytest.fixture
def pptx_group_with_table(tmp_path):
    """A table nested inside a group shape, to regression-test that
    group-nested tables are collected into tables/*.csv+*.html, not
    just rendered into document.md and then discarded."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    table_shape = slide.shapes.add_table(2, 2, Inches(0), Inches(0), Inches(3), Inches(2))
    table = table_shape.table
    table.cell(0, 0).text = "GroupCell"
    table.cell(0, 1).text = "X"
    table.cell(1, 0).text = "Y"
    table.cell(1, 1).text = "Z"
    slide.shapes._spTree.remove(table_shape._element)
    group = slide.shapes.add_group_shape([])
    group._element.append(table_shape._element)
    path = str(tmp_path / "grouped_table.pptx")
    prs.save(path)
    return path


@pytest.fixture
def pptx_with_smartart_and_ole(tmp_path):
    """Fabricates the relationship parts a real SmartArt diagram / OLE
    object would leave behind (diagramData / oleObject relationship
    types on a slide), without needing an actual PowerPoint-authored
    SmartArt object - the detector reads container relationships, not
    diagram content, so this is a faithful regression fixture."""
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Slide with a diagram and an embedded object"
    path = str(tmp_path / "special_parts.pptx")
    prs.save(path)

    import shutil
    src = path + ".src"
    shutil.move(path, src)
    with zipfile.ZipFile(src) as zin:
        original_rels = zin.read("ppt/slides/_rels/slide1.xml.rels")
        names = zin.namelist()
        contents = {n: zin.read(n) for n in names}

    tree = ET.fromstring(original_rels)
    ns = "http://schemas.openxmlformats.org/package/2006/relationships"
    ET.SubElement(tree, f"{{{ns}}}Relationship", {
        "Id": "rId100",
        "Type": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/diagramData",
        "Target": "../diagrams/data1.xml",
    })
    ET.SubElement(tree, f"{{{ns}}}Relationship", {
        "Id": "rId101",
        "Type": "http://schemas.openxmlformats.org/officeDocument/2006/relationships/oleObject",
        "Target": "../embeddings/oleObject1.bin",
    })
    contents["ppt/slides/_rels/slide1.xml.rels"] = ET.tostring(tree, encoding="UTF-8", xml_declaration=True)

    with zipfile.ZipFile(path, "w", zipfile.ZIP_DEFLATED) as zout:
        for name, data in contents.items():
            zout.writestr(name, data)
    os.remove(src)
    return path


@pytest.fixture
def pdf_column_aligned_scanned(tmp_path, monkeypatch):
    """Simulates the OCR box output for a page whose text is genuinely
    column-aligned (repeated x-positions across several lines) but with
    too few rows/columns to clear _cluster_into_table's stricter
    thresholds - the exact case TABLE_STRUCTURE_UNVERIFIED should catch,
    as opposed to a page with no alignment at all."""
    from reportlab.pdfgen import canvas
    path = str(tmp_path / "scanned_tableish.pdf")
    c = canvas.Canvas(path)
    c.showPage()
    c.save()

    import pdf_converter

    def fake_convert_scanned(pdf_path, doc, page_indices):
        # two lines, two aligned columns each - below MIN_ROWS_FOR_TABLE
        # (3) so _cluster_into_table returns None, but clearly aligned.
        boxes = [
            ([[10, 10]], "Name", 0.9),
            ([[200, 10]], "Amount", 0.9),
            ([[10, 30]], "Widget", 0.9),
            ([[200, 30]], "42", 0.9),
        ]
        elements = [{"id": "page-0001", "type": "page", "page": 1,
                     "content": "Name Amount Widget 42", "engine": "rapidocr",
                     "confidence": 0.9, "source_locator": {"page": 1}}]
        likelihood = pdf_converter._estimate_table_likelihood(boxes)
        report = {
            "status": "passed", "engine": "rapidocr_onnxruntime", "page_count": 1,
            "ocr_used": True, "ocr_avg_confidence": 0.9,
            "ocr_low_confidence_pages": [], "glued_word_pages": [],
            "tesseract_fallback_pages": [], "engine_per_page": {"1": "rapidocr"},
            "table_regions_detected": 0, "table_likelihood": likelihood,
            "table_structure_confidence": "low (no column-aligned table pattern found)",
        }
        return {"markdown": "Name Amount\n\nWidget 42", "report": report,
                "_per_page_md": ["Name Amount\n\nWidget 42"],
                "elements": elements, "tables": []}

    monkeypatch.setattr(pdf_converter, "_convert_scanned", fake_convert_scanned)
    return path


@pytest.fixture
def pdf_plain_scanned_no_table(tmp_path, monkeypatch):
    """A scanned page with no column alignment at all - the case
    TABLE_STRUCTURE_UNVERIFIED must NOT fire for (regression for the
    v1.5 blanket-warning bug)."""
    from reportlab.pdfgen import canvas
    path = str(tmp_path / "scanned_plain.pdf")
    c = canvas.Canvas(path)
    c.showPage()
    c.save()

    import pdf_converter

    def fake_convert_scanned(pdf_path, doc, page_indices):
        elements = [{"id": "page-0001", "type": "page", "page": 1,
                     "content": "just a normal sentence of prose text here",
                     "engine": "rapidocr", "confidence": 0.9,
                     "source_locator": {"page": 1}}]
        report = {
            "status": "passed", "engine": "rapidocr_onnxruntime", "page_count": 1,
            "ocr_used": True, "ocr_avg_confidence": 0.9,
            "ocr_low_confidence_pages": [], "glued_word_pages": [],
            "tesseract_fallback_pages": [], "engine_per_page": {"1": "rapidocr"},
            "table_regions_detected": 0, "table_likelihood": 0.0,
            "table_structure_confidence": "low (no column-aligned table pattern found)",
        }
        return {"markdown": "just a normal sentence of prose text here", "report": report,
                "_per_page_md": ["just a normal sentence of prose text here"],
                "elements": elements, "tables": []}

    monkeypatch.setattr(pdf_converter, "_convert_scanned", fake_convert_scanned)
    return path
