from __future__ import annotations

import json
import os
import subprocess
import time
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

ROOT = Path(__file__).resolve().parent
FIX = ROOT / "fixtures"
OUT = ROOT / "results"
FIX.mkdir(parents=True, exist_ok=True)
OUT.mkdir(parents=True, exist_ok=True)


def make_fixtures() -> list[Path]:
    files: list[Path] = []

    txt = FIX / "unicode.txt"
    txt.write_text(
        "繁體中文測試：鋼管、軸承、價格級距。\n"
        "ASCII: ERW PIPE OD 219.1 mm x WT 8.18 mm\n"
        "Symbols: ± ≤ ≥ μm °C → ← ✓\n",
        encoding="utf-8",
    )
    files.append(txt)

    csv = FIX / "merged_semantics_proxy.csv"
    csv.write_text(
        "category,range,stock_tick,bond_tick\n"
        "listed,under_5,0.01,0.01\n"
        "listed,5_to_10,0.01,0.05\n"
        "listed,10_to_50,0.05,0.05\n",
        encoding="utf-8",
    )
    files.append(csv)

    docx = FIX / "complex.docx"
    d = Document()
    d.add_heading("DOCX 結構測試", level=1)
    d.add_paragraph("段落含繁體中文、ERW、219.1 mm、±0.10 mm。")
    t = d.add_table(rows=3, cols=4)
    t.cell(0, 0).text = "分類"
    t.cell(0, 1).text = "價格範圍"
    t.cell(0, 2).text = "股票"
    t.cell(0, 3).text = "債券"
    t.cell(1, 0).text = "上市"
    t.cell(1, 1).text = "未滿5元"
    t.cell(1, 2).text = "0.01"
    t.cell(1, 3).text = "0.01"
    t.cell(2, 0).text = "上市"
    t.cell(2, 1).text = "5至10元"
    t.cell(2, 2).text = "0.01"
    t.cell(2, 3).text = "0.05"
    t.cell(1, 0).merge(t.cell(2, 0))
    d.add_heading("第二節", level=2)
    d.add_paragraph("清單：")
    for item in ["第一項", "第二項", "第三項"]:
        d.add_paragraph(item, style="List Bullet")
    d.save(docx)
    files.append(docx)

    xlsx = FIX / "complex.xlsx"
    wb = Workbook()
    ws = wb.active
    ws.title = "TickSize"
    ws.merge_cells("A1:A3")
    ws["A1"] = "上市"
    ws.append([None, "價格範圍", "股票", "債券"])
    ws.append([None, "未滿5元", 0.01, 0.01])
    ws.append(["上市", "5至10元", 0.01, 0.05])
    ws["E1"] = "公式"
    ws["E2"] = "=C3+D3"
    ws.sheet_state = "visible"
    hidden = wb.create_sheet("HiddenSheet")
    hidden["A1"] = "HIDDEN_SENTINEL"
    hidden.sheet_state = "hidden"
    wb.save(xlsx)
    files.append(xlsx)

    pptx = FIX / "complex.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "PPTX 結構測試"
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(5), Inches(1))
    box.text_frame.text = "繁體中文、ERW、219.1 mm、±0.10 mm"
    table = slide.shapes.add_table(3, 4, Inches(0.7), Inches(2.4), Inches(8), Inches(2)).table
    vals = [
        ["分類", "價格範圍", "股票", "債券"],
        ["上市", "未滿5元", "0.01", "0.01"],
        ["上市", "5至10元", "0.01", "0.05"],
    ]
    for r, row in enumerate(vals):
        for c, val in enumerate(row):
            table.cell(r, c).text = val
    table.cell(1, 0).merge(table.cell(2, 0))
    notes = slide.notes_slide.notes_text_frame
    notes.text = "SPEAKER_NOTE_SENTINEL"
    prs.save(pptx)
    files.append(pptx)

    return files


def curl(args: list[str], body_path: Path, headers_path: Path) -> dict:
    cmd = ["curl", "-sS", "-L", "--max-time", "60", "-D", str(headers_path), "-o", str(body_path), "-w", "%{http_code}|%{time_total}|%{content_type}"] + args
    started = time.time()
    p = subprocess.run(cmd, capture_output=True, text=True)
    elapsed = time.time() - started
    meta = {
        "command": cmd,
        "returncode": p.returncode,
        "curl_metrics": p.stdout.strip(),
        "stderr": p.stderr.strip(),
        "wall_seconds": round(elapsed, 3),
        "body_bytes": body_path.stat().st_size if body_path.exists() else 0,
    }
    return meta


def main() -> None:
    fixtures = make_fixtures()
    summary: list[dict] = []

    for f in fixtures:
        case = f"file_{f.suffix.lstrip('.')}"
        body = OUT / f"{case}.body.txt"
        hdr = OUT / f"{case}.headers.txt"
        meta = curl(["-X", "POST", "https://markdown.new/convert", "-F", f"file=@{f}"], body, hdr)
        meta.update({"case": case, "input": str(f.relative_to(ROOT)), "mode": "upload"})
        summary.append(meta)

    urls = [
        ("url_example", "https://example.com", "auto"),
        ("url_mdn_table", "https://developer.mozilla.org/en-US/docs/Web/HTML/Reference/Elements/table", "auto"),
        ("url_twse", "https://www.twse.com.tw/zh/products/system/trading.html", "auto"),
        ("url_twse_browser", "https://www.twse.com.tw/zh/products/system/trading.html", "browser"),
        ("url_cloudflare_blog", "https://blog.cloudflare.com/markdown-for-agents/", "auto"),
        ("url_404", "https://example.com/definitely-not-found-ufm-benchmark", "auto"),
    ]
    for case, url, method in urls:
        body = OUT / f"{case}.body.txt"
        hdr = OUT / f"{case}.headers.txt"
        payload = json.dumps({"url": url, "method": method}, ensure_ascii=False)
        meta = curl(["-X", "POST", "https://markdown.new/", "-H", "Content-Type: application/json", "--data", payload], body, hdr)
        meta.update({"case": case, "input": url, "mode": method})
        summary.append(meta)

    (OUT / "summary.json").write_text(json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8")

    report = ["# markdown.new black-box benchmark", "", "| case | return | curl metrics | bytes |", "|---|---:|---|---:|"]
    for x in summary:
        report.append(f"| {x['case']} | {x['returncode']} | `{x['curl_metrics']}` | {x['body_bytes']} |")
    (OUT / "REPORT.md").write_text("\n".join(report) + "\n", encoding="utf-8")


if __name__ == "__main__":
    main()
